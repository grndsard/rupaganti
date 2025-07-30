import os
import sqlite3
import zipfile
import threading
import time
import logging
import uuid
import shutil
import asyncio
import concurrent.futures
from datetime import datetime, timedelta
from io import BytesIO
from PIL import Image
import telebot
from telebot import types
import mimetypes
import subprocess
import base64
import secrets
import tempfile
import platform

# Import optimized encryption libraries
try:
    from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
    from cryptography.hazmat.backends import default_backend
    from cryptography.hazmat.primitives import padding
    HAS_AES = True
except ImportError:
    from cryptography.fernet import Fernet
    HAS_AES = False

# Try to import libraries for better conversions
try:
    import importlib.util
    if importlib.util.find_spec("pdf2docx"):
        print("pdf2docx library found - enhanced PDF to Word conversion available")
    else:
        print("Warning: pdf2docx not available. PDF to Word conversion may be limited.")
        print("Install with: pip install pdf2docx")
    
    if importlib.util.find_spec("docx2pdf"):
        print("docx2pdf library found - enhanced Word to PDF conversion available")
    else:
        print("Warning: docx2pdf not available. Word to PDF conversion may be limited.")
        print("Install with: pip install docx2pdf")
except ImportError:
    print("Warning: Some conversion libraries may not be available.")
    print("For best results, install: pip install pdf2docx docx2pdf")

# Try to import PDF merger library
try:
    from PyPDF2 import PdfMerger
    HAS_PDF_MERGER = True
    print("PyPDF2 library found - PDF merging available")
except ImportError:
    try:
        from pypdf import PdfMerger
        HAS_PDF_MERGER = True
        print("pypdf library found - PDF merging available")
    except ImportError:
        HAS_PDF_MERGER = False
        print("Warning: PDF merger not available. Install with: pip install PyPDF2")

# Check if we're on Windows
IS_WINDOWS = platform.system() == 'Windows'

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('bot.log', encoding='utf-8'),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger(__name__)

# Encryption setup - using AES-256 with hardware acceleration when available
if HAS_AES:
    # Generate a secure AES-256 key
    ENCRYPTION_KEY = secrets.token_bytes(32)  # 256 bits
    # Generate a secure IV (Initialization Vector)
    ENCRYPTION_IV = secrets.token_bytes(16)   # 128 bits
    # Use hardware acceleration if available
    backend = default_backend()
else:
    # Fallback to Fernet if AES is not available
    ENCRYPTION_KEY = Fernet.generate_key()
    cipher_suite = Fernet(ENCRYPTION_KEY)

# Thread pool for encryption/decryption operations
# Use more workers on multi-core systems for better performance
encryption_pool = concurrent.futures.ThreadPoolExecutor(max_workers=min(os.cpu_count() or 4, 8))

# Security settings
SECURE_DELETE_PASSES = 1  # Single pass is sufficient with modern storage
FILE_RETENTION_MINUTES = 15  # Maximum time to keep files in database
TEMP_FILE_RETENTION_MINUTES = 5  # Maximum time to keep temporary files
SESSION_TIMEOUT_SECONDS = 120  # 2-minute countdown timer for security
MIN_COMPRESSION_TARGET = 0.5  # Target at least 50% file size reduction

# Security enhancements
ALLOWED_FILE_TYPES = {'pdf', 'docx', 'doc', 'jpg', 'jpeg', 'png', 'webp', 'mp3', 'mp4'}
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB limit
RATE_LIMIT_REQUESTS = 10  # Max requests per minute per user
RATE_LIMIT_WINDOW = 60  # Time window in seconds

# Rate limiting storage
user_request_counts = {}
blocked_users = set()

# Replace with your valid Telegram bot token
BOT_TOKEN = "****"

# Configure telebot with proper request settings
telebot.apihelper.RETRY_ON_ERROR = True
telebot.apihelper.CONNECT_TIMEOUT = 10

# Initialize bot with token
bot = telebot.TeleBot(BOT_TOKEN)

# Store active sessions with their timers
active_sessions = {}

# Store user activity timestamps
user_activity = {}

# Store PDF merge sessions with enhanced batch support
pdf_merge_sessions = {}



# Store user service selections
user_services = {}

# Language translations
LANG = {
    'en': {
        'welcome': "🎉 **Welcome to RupaGanti** by Grands!\n\n🚀 Your secure file processing assistant with comprehensive tools.\n\n🛠️ **Available Services:**\n\n📄 **PDF Tools**\n• Merge multiple PDFs into one\n• Compress PDF files\n• Convert PDF to Word\n\n📸 **Image Tools**\n• Convert between JPG, PNG, WebP\n• Compress images to reduce size\n• Optimize image quality\n\n🎵 **Media Tools**\n• Convert video to MP4\n• Extract audio from videos\n• Convert audio to MP3\n\n🗜️ **Compression Tools**\n• Create ZIP archives\n• Compress any file type\n• Reduce file sizes\n\n📱 **Mobile & Desktop Optimized**\n\n🔐 **Security Features:**\n• AES-256 encryption\n• Auto-delete after processing\n• No data stored permanently\n• Local processing only\n\n👇 **Choose a service category to get started:**",
        'first_welcome': "👋 Welcome to RupaGanti by Grands!\n\nI can help you convert or compress your files safely 🔐✨\n\nTap the button below to get started 👇",
        'start_button': "🔁 Start",
        'inactivity_reminder': "👀 Are you still there?\nLet me know if you still need help!",
        'inactivity_close': "⏳ Session closed due to inactivity.\n\nIf you want to use the bot again, just tap the button below 👇",
        'convert_jpg': '📷 Convert to JPG',
        'convert_png': '🖼️ Convert to PNG',
        'convert_webp': '🌐 Convert to WebP',
        'compress_img': '🗜️ Compress Image',
        'compress_pdf': '🗜️ Compress PDF',
        'extract_mp3': '🎵 Extract Audio (MP3)',
        'zip_file': '📦 Compress to ZIP',
        'done': '✅ Done!',
        'complete': '🎁 This is what you\'ve been waiting for!',
        'help_more': '💬 Need anything else? I\'m here to help — just say the word!',
        'yes_more': '✅ Yes, Process Another',
        'no_thanks': '❌ No, Thanks!',
        'goodbye': '👋 Thanks for using RupaGanti! Type /start anytime to use me again.',
        'ready_next': '📁 Ready for your next file! Just send it to me.',
        'audio_failed': '❌ Audio extraction failed, but I\'m still here to help!',
        'encrypting': '🔐 Encrypting your file securely...',
        'encryption_complete': '✅ File encrypted successfully in {:.2f} seconds.',
        'compressing': '📉 Compressing your file...',
        'cleaning': '🧹 Cleaning up for security...',
        'session_expired': '⏱️ Session expired for your security. File has been removed.',
        'start_over': '🔁 Start Over',
        'countdown': '⏳ {}:{:02d} remaining',
        'compression_result': '✅ Compression successful: {:.1f} MB → {:.1f} MB',
        'already_optimized': '⚠️ This file is already optimized and cannot be compressed further without losing quality.',
        'files_deleted': '🗑️ All files have been securely deleted to protect your data.',
        'file_received': 'received!\n\nChoose what you\'d like to do:',
        'error_processing': '❌ Sorry, there was a problem processing your file.',
        'error_upload': '❌ File upload failed. Please try uploading your file again.',
        'error_general': '❌ Something went wrong. Please try again.',
        'back_to_start': '🔁 Back to Start',
        'try_again': 'Please try uploading your file again or contact support if the problem persists.',
        'convert_to_pdf': '📄 Convert to PDF',
        'what_to_do': 'What would you like to do?',
        'file_ready': '🎉 Done! Your file is ready. Download below 👇',
        'security_reminder': '🔒 Your files are encrypted and will be deleted within 2 minutes if not processed — your data stays safe!',
        'oops_error': '⚠️ Oops! Something went wrong. Please try again or start over.',
        'converting_to_pdf': '📄 Converting to PDF...',
        'pdf_conversion_success': '✅ Word document converted to PDF successfully!',
        'pdf_conversion_failed': '❌ PDF conversion failed. Please try again with a different file.',
        'unsupported_format': '⚠️ Unsupported file format. Please upload one of the supported types:\nImages: JPG, PNG, WebP, BMP, etc.\nDocuments: DOCX, XLSX, PPTX, PDF, etc.\nAudio/Video: MP3, MP4, AVI, etc.',
        'convert_document': '📄 Convert to PDF',
        'compress_document': '🗜️ Compress to ZIP',
        'image_options': 'What would you like to do with your image?',
        'document_options': 'What would you like to do with your document?',
        'media_options': 'What would you like to do with your media file?',
        'convert_to_mp4': '🎬 Convert to MP4',
        'convert_to_mp3': '🎵 Convert to MP3',
        'cancel': '❌ Cancel',
        'combine_pdf': '📄 Combine PDFs',
        'pdf_merge_start': '📄 PDF Merge Mode\n\nSend me 2 or more PDF files to combine them into one.\n\nI\'ll show you the order and let you rearrange before merging.',
        'pdf_added': '✅ PDF added: {}\nPosition: {}\n\nSend more PDFs or tap "Merge Now" when ready.',
        'pdf_merge_ready': '📄 Ready to merge {} PDFs:\n\n{}\n\nTap buttons to reorder or merge now.',
        'merge_now': '🔗 Merge Now',
        'move_up': '⬆️ Move Up',
        'move_down': '⬇️ Move Down',
        'remove_pdf': '❌ Remove',
        'pdf_merging': '🔄 Merging PDFs...',
        'pdf_merge_success': '✅ PDFs merged successfully!',
        'pdf_merge_failed': '❌ PDF merge failed. Please try again.',
        'pdf_merge_min_files': '⚠️ You need at least 2 PDF files to merge.',
        'pdf_merge_cancelled': '❌ PDF merge cancelled.',
        'pdf_only_merge': '⚠️ Only PDF files can be added to merge session.',
        'pdf_merge_limit': '⚠️ Maximum 10 PDFs can be merged at once.',
        'pdf_file_corrupted': '❌ PDF file appears to be corrupted or invalid.',
        'pdf_merge_timeout': '⏰ PDF merge session expired. Start over to merge PDFs.',
        'cancel_merge': '❌ Cancel Merge',
        'service_menu': 'Choose a service:',
        'pdf_tools': '📄 PDF Tools',
        'image_tools': '📸 Image Tools', 
        'media_tools': '🎵 Media Tools',
        'compression_tools': '🗜️ Compression',
        'upload_pdf_merge': '📄 Upload 2+ PDF files to merge them together',
        'upload_pdf_compress': '📄 Upload a PDF file to compress it', 
        'upload_pdf_convert': '📄 Upload a PDF to convert to Word',
        'upload_pdf_split': '✂️ Upload a PDF file to split into separate pages',
        'upload_image': '📸 Upload an image file (JPG, PNG, WebP, etc.)',
        'upload_media': '🎵 Upload audio/video file (MP3, MP4, etc.)',
        'upload_document': '📄 Upload document to convert to PDF',
        'wrong_file_type': '❌ Wrong file type! Please upload the correct file type for this service.',
        'back_to_menu': '🔙 Back to Menu'
    },
    'id': {
        'welcome': "🎉 **Selamat datang di RupaGanti** by Grands!\n\n🚀 Asisten pemrosesan file aman dengan tools lengkap.\n\n🛠️ **Layanan Tersedia:**\n\n📄 **Tools PDF**\n• Gabung beberapa PDF jadi satu\n• Kompres file PDF\n• Konversi PDF ke Word\n\n📸 **Tools Gambar**\n• Konversi antara JPG, PNG, WebP\n• Kompres gambar untuk kurangi ukuran\n• Optimasi kualitas gambar\n\n🎵 **Tools Media**\n• Konversi video ke MP4\n• Ekstrak audio dari video\n• Konversi audio ke MP3\n\n🗜️ **Tools Kompresi**\n• Buat arsip ZIP\n• Kompres semua jenis file\n• Kurangi ukuran file\n\n📱 **Dioptimalkan untuk Mobile & Desktop**\n\n🔐 **Fitur Keamanan:**\n• Enkripsi AES-256\n• Hapus otomatis setelah proses\n• Tidak ada data disimpan permanen\n• Proses lokal saja\n\n👇 **Pilih kategori layanan untuk memulai:**",
        'first_welcome': '👋 Selamat datang di RupaGanti by Grands!\n\nSaya dapat membantu Anda mengkonversi atau mengompres file dengan aman 🔐✨\n\nKetuk tombol di bawah untuk memulai 👇',
        'start_button': '🔁 Mulai',
        'inactivity_reminder': '👀 Apakah Anda masih di sana?\nBeri tahu saya jika Anda masih membutuhkan bantuan!',
        'inactivity_close': '⏳ Sesi ditutup karena tidak aktif.\n\nJika Anda ingin menggunakan bot lagi, cukup ketuk tombol di bawah 👇',
        'convert_jpg': '📷 Konversi ke JPG',
        'convert_png': '🖼️ Konversi ke PNG',
        'convert_webp': '🌐 Konversi ke WebP',
        'compress_img': '🗜️ Kompres Gambar',
        'compress_pdf': '🗜️ Kompres PDF',
        'extract_mp3': '🎵 Ekstrak Audio (MP3)',
        'zip_file': '📦 Kompres ke ZIP',
        'done': '✅ Selesai!',
        'complete': '🎁 Inilah yang sudah Anda tunggu!',
        'help_more': '💬 Butuh bantuan lagi? Saya siap membantu — tinggal bilang saja!',
        'yes_more': '✅ Ya, Proses Lagi',
        'no_thanks': '❌ Tidak, Terima Kasih!',
        'goodbye': '👋 Terima kasih telah menggunakan RupaGanti! Ketik /start kapan saja untuk menggunakan saya lagi.',
        'ready_next': '📁 Siap untuk file berikutnya! Kirim saja ke saya.',
        'audio_failed': '❌ Ekstraksi audio gagal, tapi saya masih di sini untuk membantu!',
        'encrypting': '🔐 Mengenkripsi file Anda dengan aman...',
        'encryption_complete': '✅ File berhasil dienkripsi dalam {:.2f} detik.',
        'compressing': '📉 Mengompres file Anda...',
        'cleaning': '🧹 Membersihkan untuk keamanan...',
        'session_expired': '⏱️ Sesi berakhir untuk keamanan Anda. File telah dihapus.',
        'start_over': '🔁 Mulai Lagi',
        'countdown': '⏳ {}:{:02d} tersisa',
        'compression_result': '✅ Kompresi berhasil: {:.1f} MB → {:.1f} MB',
        'already_optimized': '⚠️ File ini sudah dioptimalkan dan tidak dapat dikompres lebih lanjut tanpa kehilangan kualitas.',
        'files_deleted': '🗑️ Semua file telah dihapus dengan aman untuk melindungi data Anda.',
        'file_received': 'diterima!\n\nPilih apa yang ingin Anda lakukan:',
        'error_processing': '❌ Maaf, ada masalah saat memproses file Anda.',
        'error_upload': '❌ Upload file gagal. Silakan coba upload file Anda lagi.',
        'error_general': '❌ Ada yang salah. Silakan coba lagi.',
        'back_to_start': '🔁 Kembali ke Awal',
        'try_again': 'Silakan coba upload file Anda lagi atau hubungi dukungan jika masalah berlanjut.',
        'convert_to_pdf': '📄 Konversi ke PDF',
        'what_to_do': 'Apa yang ingin Anda lakukan?',
        'file_ready': '🎉 Selesai! File Anda siap. Unduh di bawah 👇',
        'security_reminder': '🔒 File Anda dienkripsi dan akan dihapus dalam 2 menit jika tidak diproses — data Anda tetap aman!',
        'oops_error': '⚠️ Ups! Ada yang salah. Silakan coba lagi atau mulai ulang.',
        'converting_to_pdf': '📄 Mengkonversi ke PDF...',
        'pdf_conversion_success': '✅ Dokumen Word berhasil dikonversi ke PDF!',
        'pdf_conversion_failed': '❌ Konversi PDF gagal. Silakan coba lagi dengan file yang berbeda.',
        'unsupported_format': '⚠️ Format file tidak didukung. Silakan upload salah satu jenis yang didukung:\nGambar: JPG, PNG, WebP, BMP, dll.\nDokumen: DOCX, XLSX, PPTX, PDF, dll.\nAudio/Video: MP3, MP4, AVI, dll.',
        'convert_document': '📄 Konversi ke PDF',
        'compress_document': '🗜️ Kompres ke ZIP',
        'image_options': 'Apa yang ingin Anda lakukan dengan gambar Anda?',
        'document_options': 'Apa yang ingin Anda lakukan dengan dokumen Anda?',
        'media_options': 'Apa yang ingin Anda lakukan dengan file media Anda?',
        'convert_to_mp4': '🎬 Konversi ke MP4',
        'convert_to_mp3': '🎵 Konversi ke MP3',
        'cancel': '❌ Batal',
        'combine_pdf': '📄 Gabung PDF',
        'pdf_merge_start': '📄 Mode Gabung PDF\n\nKirim 2 atau lebih file PDF untuk menggabungkannya.\n\nSaya akan tunjukkan urutan dan biarkan Anda mengatur ulang sebelum menggabung.',
        'pdf_added': '✅ PDF ditambahkan: {}\nPosisi: {}\n\nKirim PDF lagi atau tap "Gabung Sekarang" jika siap.',
        'pdf_merge_ready': '📄 Siap menggabung {} PDF:\n\n{}\n\nTap tombol untuk mengatur ulang atau gabung sekarang.',
        'merge_now': '🔗 Gabung Sekarang',
        'move_up': '⬆️ Naik',
        'move_down': '⬇️ Turun',
        'remove_pdf': '❌ Hapus',
        'pdf_merging': '🔄 Menggabung PDF...',
        'pdf_merge_success': '✅ PDF berhasil digabung!',
        'pdf_merge_failed': '❌ Gagal menggabung PDF. Silakan coba lagi.',
        'pdf_merge_min_files': '⚠️ Anda perlu minimal 2 file PDF untuk digabung.',
        'pdf_merge_cancelled': '❌ Penggabungan PDF dibatalkan.',
        'pdf_only_merge': '⚠️ Hanya file PDF yang bisa ditambahkan ke sesi gabung.',
        'pdf_merge_limit': '⚠️ Maksimal 10 PDF bisa digabung sekaligus.',
        'pdf_file_corrupted': '❌ File PDF tampaknya rusak atau tidak valid.',
        'pdf_merge_timeout': '⏰ Sesi gabung PDF berakhir. Mulai ulang untuk gabung PDF.',
        'cancel_merge': '❌ Batal Gabung',
        'service_menu': 'Pilih layanan:',
        'pdf_tools': '📄 Tools PDF',
        'image_tools': '📸 Tools Gambar',
        'media_tools': '🎵 Tools Media',
        'compression_tools': '🗜️ Kompresi',
        'upload_pdf_merge': '📄 Upload 2+ file PDF untuk digabung',
        'upload_pdf_compress': '📄 Upload file PDF untuk dikompres',
        'upload_pdf_convert': '📄 Upload PDF untuk dikonversi ke Word',
        'upload_pdf_split': '✂️ Upload file PDF untuk dipisah menjadi halaman terpisah',
        'upload_image': '📸 Upload file gambar (JPG, PNG, WebP, dll.)',
        'upload_media': '🎵 Upload file audio/video (MP3, MP4, dll.)',
        'upload_document': '📄 Upload dokumen untuk dikonversi ke PDF',
        'wrong_file_type': '❌ Tipe file salah! Silakan upload tipe file yang benar untuk layanan ini.',
        'back_to_menu': '🔙 Kembali ke Menu'
    },
    'ar': {
        'welcome': "🎉 **مرحباً بك في RupaGanti** من Grands!\n\n🚀 مساعدك الآمن لمعالجة الملفات مع أدوات شاملة.\n\n🛠️ **الخدمات المتاحة:**\n\n📄 **أدوات PDF**\n• دمج عدة ملفات PDF في واحد\n• ضغط ملفات PDF\n• تحويل PDF إلى Word\n\n📸 **أدوات الصور**\n• تحويل بين JPG، PNG، WebP\n• ضغط الصور لتقليل الحجم\n• تحسين جودة الصور\n\n🎵 **أدوات الوسائط**\n• تحويل الفيديو إلى MP4\n• استخراج الصوت من الفيديو\n• تحويل الصوت إلى MP3\n\n🗜️ **أدوات الضغط**\n• إنشاء أرشيف ZIP\n• ضغط أي نوع ملف\n• تقليل أحجام الملفات\n\n📱 **محسّن للهاتف وسطح المكتب**\n\n🔐 **ميزات الأمان:**\n• تشفير AES-256\n• حذف تلقائي بعد المعالجة\n• لا يتم حفظ البيانات بشكل دائم\n• معالجة محلية فقط\n\n👇 **اختر فئة خدمة للبدء:**",
        'first_welcome': '👋 مرحباً بك في RupaGanti by Grands!\n\nيمكنني مساعدتك في تحويل أو ضغط ملفاتك بأمان 🔐✨\n\nاضغط على الزر أدناه للبدء 👇',
        'start_button': '🔁 ابدأ',
        'inactivity_reminder': '👀 هل ما زلت هناك؟\nأخبرني إذا كنت لا تزال بحاجة إلى مساعدة!',
        'inactivity_close': '⏳ تم إغلاق الجلسة بسبب عدم النشاط.\n\nإذا كنت تريد استخدام البوت مرة أخرى، فقط اضغط على الزر أدناه 👇',
        'convert_jpg': '📷 تحويل إلى JPG',
        'convert_png': '🖼️ تحويل إلى PNG',
        'convert_webp': '🌐 تحويل إلى WebP',
        'compress_img': '🗜️ ضغط الصورة',
        'compress_pdf': '🗜️ ضغط PDF',
        'extract_mp3': '🎵 استخراج الصوت (MP3)',
        'zip_file': '📦 ضغط إلى ZIP',
        'done': '✅ تم!',
        'complete': '🎁 هذا ما كنت تنتظره!',
        'help_more': '💬 تحتاج شيئاً آخر؟ أنا هنا للمساعدة — فقط قل كلمة!',
        'yes_more': '✅ نعم، معالجة أخرى',
        'no_thanks': '❌ لا، شكراً!',
        'goodbye': '👋 شكراً لاستخدام RupaGanti! اكتب /start في أي وقت لاستخدامي مرة أخرى.',
        'ready_next': '📁 جاهز لملفك التالي! أرسله لي فقط.',
        'audio_failed': '❌ فشل استخراج الصوت، لكنني ما زلت هنا للمساعدة!',
        'encrypting': '🔐 تشفير ملفك بأمان...',
        'encryption_complete': '✅ تم تشفير الملف بنجاح في {:.2f} ثانية.',
        'compressing': '📉 ضغط ملفك...',
        'cleaning': '🧹 التنظيف للأمان...',
        'session_expired': '⏱️ انتهت الجلسة لأمانك. تم حذف الملف.',
        'start_over': '🔁 ابدأ من جديد',
        'countdown': '⏳ {}:{:02d} متبقي',
        'compression_result': '✅ نجح الضغط: {:.1f} ميجابايت → {:.1f} ميجابايت',
        'already_optimized': '⚠️ هذا الملف محسّن بالفعل ولا يمكن ضغطه أكثر دون فقدان الجودة.',
        'files_deleted': '🗑️ تم حذف جميع الملفات بأمان لحماية بياناتك.',
        'file_received': 'تم الاستلام!\n\nاختر ما تريد فعله:',
        'error_processing': '❌ عذراً، هناك مشكلة في معالجة ملفك.',
        'error_upload': '❌ فشل رفع الملف. يرجى المحاولة مرة أخرى.',
        'error_general': '❌ حدث خطأ. يرجى المحاولة مرة أخرى.',
        'back_to_start': '🔁 العودة للبداية',
        'try_again': 'يرجى المحاولة مرة أخرى أو الاتصال بالدعم إذا استمرت المشكلة.',
        'convert_to_pdf': '📄 تحويل إلى PDF',
        'what_to_do': 'ماذا تريد أن تفعل؟',
        'file_ready': '🎉 تم! ملفك جاهز. حمّل من أدناه 👇',
        'security_reminder': '🔒 ملفاتك مشفرة وسيتم حذفها خلال دقيقتين إذا لم تتم معالجتها — بياناتك آمنة!',
        'oops_error': '⚠️ عفواً! حدث خطأ. يرجى المحاولة مرة أخرى أو البدء من جديد.',
        'converting_to_pdf': '📄 تحويل إلى PDF...',
        'pdf_conversion_success': '✅ تم تحويل مستند Word إلى PDF بنجاح!',
        'pdf_conversion_failed': '❌ فشل تحويل PDF. يرجى المحاولة مع ملف مختلف.',
        'unsupported_format': '⚠️ صيغة ملف غير مدعومة. يرجى رفع أحد الأنواع المدعومة:\nالصور: JPG, PNG, WebP, BMP, إلخ.\nالمستندات: DOCX, XLSX, PPTX, PDF, إلخ.\nالصوت/الفيديو: MP3, MP4, AVI, إلخ.',
        'convert_document': '📄 تحويل إلى PDF',
        'compress_document': '🗜️ ضغط إلى ZIP',
        'image_options': 'ماذا تريد أن تفعل بصورتك؟',
        'document_options': 'ماذا تريد أن تفعل بمستندك؟',
        'media_options': 'ماذا تريد أن تفعل بملف الوسائط؟',
        'convert_to_mp4': '🎬 تحويل إلى MP4',
        'convert_to_mp3': '🎵 تحويل إلى MP3',
        'cancel': '❌ إلغاء',
        'combine_pdf': '📄 دمج PDF',
        'pdf_merge_start': '📄 وضع دمج PDF\n\nأرسل لي ملفين PDF أو أكثر لدمجهما في ملف واحد.\n\nسأعرض عليك الترتيب وأتيح لك إعادة الترتيب قبل الدمج.',
        'pdf_added': '✅ تم إضافة PDF: {}\nالموضع: {}\n\nأرسل المزيد من ملفات PDF أو اضغط "ادمج الآن" عند الاستعداد.',
        'pdf_merge_ready': '📄 جاهز لدمج {} ملفات PDF:\n\n{}\n\nاضغط الأزرار لإعادة الترتيب أو ادمج الآن.',
        'merge_now': '🔗 ادمج الآن',
        'move_up': '⬆️ تحريك لأعلى',
        'move_down': '⬇️ تحريك لأسفل',
        'remove_pdf': '❌ إزالة',
        'pdf_merging': '🔄 دمج ملفات PDF...',
        'pdf_merge_success': '✅ تم دمج ملفات PDF بنجاح!',
        'pdf_merge_failed': '❌ فشل دمج PDF. يرجى المحاولة مرة أخرى.',
        'pdf_merge_min_files': '⚠️ تحتاج إلى ملفين PDF على الأقل للدمج.',
        'pdf_merge_cancelled': '❌ تم إلغاء دمج PDF.',
        'pdf_only_merge': '⚠️ يمكن إضافة ملفات PDF فقط لجلسة الدمج.',
        'pdf_merge_limit': '⚠️ يمكن دمج 10 ملفات PDF كحد أقصى في المرة الواحدة.',
        'pdf_file_corrupted': '❌ يبدو أن ملف PDF تالف أو غير صالح.',
        'pdf_merge_timeout': '⏰ انتهت جلسة دمج PDF. ابدأ من جديد لدمج ملفات PDF.',
        'cancel_merge': '❌ إلغاء الدمج'
    },
    'jv': {
        'welcome': "🎉 **Sugeng rawuh ing RupaGanti** saka Grands!\n\n🚀 Asisten pangolahan file aman karo tools lengkap.\n\n🛠️ **Layanan sing Ana:**\n\n📄 **Tools PDF**\n• Gabung pirang-pirang PDF dadi siji\n• Kompres file PDF\n• Konversi PDF dadi Word\n\n📸 **Tools Gambar**\n• Konversi antarane JPG, PNG, WebP\n• Kompres gambar kanggo ngurangi ukuran\n• Optimasi kualitas gambar\n\n🎵 **Tools Media**\n• Konversi video dadi MP4\n• Ekstrak audio saka video\n• Konversi audio dadi MP3\n\n🗜️ **Tools Kompresi**\n• Gawe arsip ZIP\n• Kompres kabeh jinis file\n• Ngurangi ukuran file\n\n📱 **Dioptimalake kanggo Mobile & Desktop**\n\n🔐 **Fitur Keamanan:**\n• Enkripsi AES-256\n• Busak otomatis sawise proses\n• Ora ana data disimpen permanen\n• Proses lokal wae\n\n👇 **Pilih kategori layanan kanggo miwiti:**",
        'first_welcome': '👋 Sugeng rawuh ing RupaGanti by Grands!\n\nAku bisa ngewangi sampeyan ngowahi utawa ngompres file kanthi aman 🔐✨\n\nKlik tombol ing ngisor kanggo miwiti 👇',
        'start_button': '🔁 Miwiti',
        'inactivity_reminder': '👀 Sampeyan isih ana?\nKandha yen sampeyan isih butuh bantuan!',
        'inactivity_close': '⏳ Sesi ditutup amarga ora aktif.\n\nYen sampeyan arep nggunakake bot maneh, klik tombol ing ngisor 👇',
        'convert_jpg': '📷 Owahi dadi JPG',
        'convert_png': '🖼️ Owahi dadi PNG',
        'convert_webp': '🌐 Owahi dadi WebP',
        'compress_img': '🗜️ Kompres Gambar',
        'compress_pdf': '🗜️ Kompres PDF',
        'extract_mp3': '🎵 Jupuk Audio (MP3)',
        'zip_file': '📦 Kompres dadi ZIP',
        'done': '✅ Rampung!',
        'complete': '🎁 Iki sing wis kokenteni!',
        'help_more': '💬 Butuh apa maneh? Aku kene kanggo ngewangi — cukup kandha wae!',
        'yes_more': '✅ Ya, Proses Maneh',
        'no_thanks': '❌ Ora, Matur Nuwun!',
        'goodbye': '👋 Matur nuwun wis nggunakake RupaGanti! Ketik /start kapan wae kanggo nggunakake aku maneh.',
        'ready_next': '📁 Siap kanggo file sabanjure! Kirimno wae marang aku.',
        'audio_failed': '❌ Ekstraksi audio gagal, nanging aku isih kene kanggo ngewangi!',
        'encrypting': '🔐 Ngenkripsi file sampeyan kanthi aman...',
        'encryption_complete': '✅ File kasil dienkripsi sajrone {:.2f} detik.',
        'compressing': '📉 Ngompres file sampeyan...',
        'cleaning': '🧹 Ngresiki kanggo keamanan...',
        'session_expired': '⏱️ Sesi rampung kanggo keamanan sampeyan. File wis dibusak.',
        'start_over': '🔁 Miwiti Maneh',
        'countdown': '⏳ {}:{:02d} isih ana',
        'compression_result': '✅ Kompresi kasil: {:.1f} MB → {:.1f} MB',
        'already_optimized': '⚠️ File iki wis dioptimalake lan ora bisa dikompres maneh tanpa kelangan kualitas.',
        'files_deleted': '🗑️ Kabeh file wis dibusak kanthi aman kanggo nglindhungi data sampeyan.',
        'file_received': 'ditampa!\n\nPilih apa sing arep koklakoni:',
        'error_processing': '❌ Nuwun sewu, ana masalah nalika ngolah file sampeyan.',
        'error_upload': '❌ Upload file gagal. Coba upload file sampeyan maneh.',
        'error_general': '❌ Ana sing salah. Coba maneh.',
        'back_to_start': '🔁 Bali menyang Wiwitan',
        'try_again': 'Coba upload file sampeyan maneh utawa hubungi dhukungan yen masalah terus.',
        'convert_to_pdf': '📄 Owahi dadi PDF',
        'what_to_do': 'Apa sing arep koklakoni?',
        'file_ready': '🎉 Rampung! File sampeyan siap. Download ing ngisor 👇',
        'security_reminder': '🔒 File sampeyan dienkripsi lan bakal dibusak sajrone 2 menit yen ora diproses — data sampeyan tetep aman!',
        'oops_error': '⚠️ Aduh! Ana sing salah. Coba maneh utawa miwiti maneh.',
        'converting_to_pdf': '📄 Ngowahi dadi PDF...',
        'pdf_conversion_success': '✅ Dokumen Word kasil diowahi dadi PDF!',
        'pdf_conversion_failed': '❌ Konversi PDF gagal. Coba maneh karo file liyane.',
        'unsupported_format': '⚠️ Format file ora didukung. Upload salah siji jinis sing didukung:\nGambar: JPG, PNG, WebP, BMP, lsp.\nDokumen: DOCX, XLSX, PPTX, PDF, lsp.\nAudio/Video: MP3, MP4, AVI, lsp.',
        'convert_document': '📄 Owahi dadi PDF',
        'compress_document': '🗜️ Kompres dadi ZIP',
        'image_options': 'Apa sing arep koklakoni karo gambar sampeyan?',
        'document_options': 'Apa sing arep koklakonke karo dokumen sampeyan?',
        'media_options': 'Apa sing arep koklakoni karo file media sampeyan?',
        'convert_to_mp4': '🎬 Owahi dadi MP4',
        'convert_to_mp3': '🎵 Owahi dadi MP3',
        'cancel': '❌ Batal',
        'combine_pdf': '📄 Gabung PDF',
        'pdf_merge_start': '📄 Mode Gabung PDF\n\nKirimno 2 utawa luwih file PDF kanggo nggabungake dadi siji.\n\nAku bakal nuduhake urutan lan ngidini sampeyan ngatur maneh sadurunge nggabung.',
        'pdf_added': '✅ PDF ditambahake: {}\nPosisi: {}\n\nKirim PDF liyane utawa pencet "Gabung Saiki" yen wis siap.',
        'pdf_merge_ready': '📄 Siap nggabung {} PDF:\n\n{}\n\nPencet tombol kanggo ngatur maneh utawa gabung saiki.',
        'merge_now': '🔗 Gabung Saiki',
        'move_up': '⬆️ Munggah',
        'move_down': '⬇️ Mudhun',
        'remove_pdf': '❌ Busak',
        'pdf_merging': '🔄 Nggabung PDF...',
        'pdf_merge_success': '✅ PDF kasil digabung!',
        'pdf_merge_failed': '❌ Gagal nggabung PDF. Coba maneh.',
        'pdf_merge_min_files': '⚠️ Sampeyan butuh minimal 2 file PDF kanggo digabung.',
        'pdf_merge_cancelled': '❌ Penggabungan PDF dibatalake.',
        'pdf_only_merge': '⚠️ Mung file PDF sing bisa ditambahake ing sesi gabung.',
        'pdf_merge_limit': '⚠️ Maksimal 10 PDF bisa digabung sekaligus.',
        'pdf_file_corrupted': '❌ File PDF katon rusak utawa ora valid.',
        'pdf_merge_timeout': '⏰ Sesi gabung PDF rampung. Miwiti maneh kanggo gabung PDF.',
        'cancel_merge': '❌ Batal Gabung',
        'upload_pdf_split': '✂️ Upload file PDF kanggo dipisah dadi kaca-kaca terpisah'
    }
}

def get_user_lang(lang_code):
    """
    Menentukan bahasa pengguna berdasarkan kode bahasa Telegram.
    
    Parameter:
        lang_code (str): Kode bahasa dari Telegram (contoh: 'id', 'en', 'ar')
    
    Return:
        str: Kode bahasa yang didukung ('id', 'ar', 'jv', atau 'en' sebagai default)
    
    Catatan:
        - Bahasa Indonesia: 'id'
        - Bahasa Arab: 'ar' 
        - Bahasa Jawa: 'jv'
        - Bahasa Inggris: 'en' (default)
    """
    if lang_code and lang_code.startswith('id'):
        return 'id'
    elif lang_code and lang_code.startswith('ar'):
        return 'ar'
    elif lang_code and lang_code.startswith('jv'):
        return 'jv'
    return 'en'

os.makedirs("files", exist_ok=True)
os.makedirs("temp", exist_ok=True)

def init_db():
    """
    Menginisialisasi database SQLite untuk menyimpan informasi file.
    
    Parameter:
        Tidak ada
    
    Return:
        Tidak ada
    
    Catatan:
        - Membuat tabel 'files' jika belum ada
        - Tabel berisi: id, user_id, file_id, file_name, file_path, created_at
        - Menggunakan timeout 10 detik untuk koneksi database
        - Akan raise exception jika inisialisasi gagal
    """
    try:
        conn = sqlite3.connect('files.db', timeout=10.0)
        conn.execute('''CREATE TABLE IF NOT EXISTS files
                        (id INTEGER PRIMARY KEY, user_id INTEGER, file_id TEXT, 
                         file_name TEXT, file_path TEXT, created_at TIMESTAMP)''')
        conn.commit()
        conn.close()
        logger.info("Database initialized successfully")
    except Exception as e:
        logger.error(f"Database initialization failed: {str(e)}")
        raise

init_db()

def get_file_type(filename):
    """
    Menentukan jenis file berdasarkan ekstensi nama file.
    
    Parameter:
        filename (str): Nama file dengan ekstensi
    
    Return:
        tuple: (kategori_file, ekstensi)
               - kategori_file: 'image', 'document', 'video', 'audio', atau 'unsupported'
               - ekstensi: ekstensi file dalam huruf kecil
    
    Catatan:
        - Mendukung berbagai format file populer
        - Gambar: jpg, jpeg, png, webp, bmp, gif, tiff
        - Dokumen: pdf, doc, docx, txt, rtf, xlsx, xls, pptx, ppt
        - Video: mp4, avi, mov, mkv, wmv
        - Audio: mp3, wav, flac, aac, m4a
    """
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    file_types = {
        'image': ['jpg', 'jpeg', 'png', 'webp', 'bmp', 'gif', 'tiff'],
        'document': ['pdf', 'doc', 'docx', 'txt', 'rtf', 'xlsx', 'xls', 'pptx', 'ppt'],
        'video': ['mp4', 'avi', 'mov', 'mkv', 'wmv'],
        'audio': ['mp3', 'wav', 'flac', 'aac', 'm4a']
    }
    
    for category, extensions in file_types.items():
        if ext in extensions:
            return category, ext
    return 'unsupported', ext

def is_supported_file(filename):
    """
    Memeriksa apakah file dalam format yang didukung.
    
    Parameter:
        filename (str): Nama file yang akan diperiksa
    
    Return:
        bool: True jika file didukung, False jika tidak
    
    Catatan:
        - Menggunakan fungsi get_file_type() untuk menentukan jenis file
        - File didukung jika bukan kategori 'unsupported'
    """
    file_type, ext = get_file_type(filename)
    return file_type != 'unsupported'

def encrypt_file_aes(file_data):
    """
    Mengenkripsi data file menggunakan enkripsi AES-256 dengan akselerasi hardware.
    
    Parameter:
        file_data (bytes): Data file yang akan dienkripsi
    
    Return:
        bytes: Data terenkripsi dengan IV di awal
    
    Catatan:
        - Menggunakan AES-256 dalam mode CBC
        - Untuk file besar (>10MB), memproses dalam chunk 1MB untuk efisiensi memori
        - Untuk file kecil, menggunakan pendekatan sederhana
        - IV (Initialization Vector) ditambahkan di awal data terenkripsi
        - Menggunakan PKCS7 padding
        - Akan raise exception jika enkripsi gagal
    """
    try:
        # For large files, use a more efficient approach with less memory overhead
        if len(file_data) > 10 * 1024 * 1024:  # 10MB
            # Process in chunks for large files
            chunk_size = 1024 * 1024  # 1MB chunks
            padder = padding.PKCS7(algorithms.AES.block_size).padder()
            cipher = Cipher(algorithms.AES(ENCRYPTION_KEY), modes.CBC(ENCRYPTION_IV), backend=backend)
            encryptor = cipher.encryptor()
            
            # Process all but the last chunk
            result = bytearray()
            for i in range(0, len(file_data) - chunk_size, chunk_size):
                chunk = file_data[i:i+chunk_size]
                result.extend(encryptor.update(chunk))
            
            # Process the last chunk with padding
            last_chunk = file_data[-(len(file_data) % chunk_size or chunk_size):]
            padded_last_chunk = padder.update(last_chunk) + padder.finalize()
            result.extend(encryptor.update(padded_last_chunk) + encryptor.finalize())
            
            # Prepend IV to the encrypted data for decryption later
            return ENCRYPTION_IV + bytes(result)
        else:
            # For smaller files, use the simpler approach
            padder = padding.PKCS7(algorithms.AES.block_size).padder()
            padded_data = padder.update(file_data) + padder.finalize()
            
            cipher = Cipher(algorithms.AES(ENCRYPTION_KEY), modes.CBC(ENCRYPTION_IV), backend=backend)
            encryptor = cipher.encryptor()
            
            encrypted_data = encryptor.update(padded_data) + encryptor.finalize()
            
            # Prepend IV to the encrypted data for decryption later
            return ENCRYPTION_IV + encrypted_data
    except Exception as e:
        logger.error(f"AES encryption error: {str(e)}")
        raise

def encrypt_file_fernet(file_data):
    """
    Mengenkripsi data file menggunakan enkripsi Fernet.
    
    Parameter:
        file_data (bytes): Data file yang akan dienkripsi
    
    Return:
        bytes: Data terenkripsi menggunakan Fernet
    
    Catatan:
        - Fernet adalah enkripsi simetris yang aman dan mudah digunakan
        - Digunakan sebagai fallback jika AES tidak tersedia
        - Akan raise exception jika enkripsi gagal
    """
    try:
        return cipher_suite.encrypt(file_data)
    except Exception as e:
        logger.error(f"Fernet encryption error: {str(e)}")
        raise

def encrypt_file(file_data):
    """
    Mengenkripsi data file menggunakan metode terbaik yang tersedia.
    
    Parameter:
        file_data (bytes): Data file yang akan dienkripsi
    
    Return:
        tuple: (data_terenkripsi, waktu_enkripsi)
               - data_terenkripsi (bytes): Data yang sudah dienkripsi
               - waktu_enkripsi (float): Waktu yang dibutuhkan untuk enkripsi dalam detik
    
    Catatan:
        - Prioritas: AES-256 (jika tersedia) > Fernet (fallback)
        - Mengukur waktu enkripsi untuk monitoring performa
        - Jika semua enkripsi gagal, mengembalikan data asli sebagai last resort
        - Log semua aktivitas enkripsi untuk debugging
    """
    start_time = time.time()
    try:
        if HAS_AES:
            # Use AES-256 encryption with hardware acceleration
            result = encrypt_file_aes(file_data)
        else:
            # Fallback to Fernet encryption
            result = encrypt_file_fernet(file_data)
        
        encryption_time = time.time() - start_time
        logger.info(f"File encrypted successfully in {encryption_time:.2f} seconds")
        return result, encryption_time
    except Exception as e:
        logger.error(f"Encryption error: {str(e)}")
        # If encryption fails, return original data with warning
        logger.critical("Encryption completely failed, returning original data")
        return file_data, 0  # Last resort
        
def get_file_size_mb(data):
    """
    Menghitung ukuran file dalam MB (Megabyte).
    
    Parameter:
        data (bytes): Data file yang akan dihitung ukurannya
    
    Return:
        float: Ukuran file dalam MB
    
    Catatan:
        - Konversi: bytes → KB (÷1024) → MB (÷1024)
        - Hasil dalam format desimal untuk presisi
    """
    return len(data) / (1024 * 1024)

def calculate_compression_ratio(original_size, compressed_size):
    """
    Menghitung rasio kompresi file.
    
    Parameter:
        original_size (float): Ukuran file asli
        compressed_size (float): Ukuran file setelah kompresi
    
    Return:
        float: Rasio kompresi (0.0 - 1.0)
               - 0.0 = tidak ada kompresi
               - 1.0 = kompresi sempurna (100%)
    
    Catatan:
        - Rumus: 1 - (ukuran_terkompresi / ukuran_asli)
        - Menangani kasus ukuran asli = 0 untuk menghindari division by zero
    """
    if original_size == 0:
        return 0
    return 1 - (compressed_size / original_size)

def decrypt_file_aes(encrypted_data):
    """
    Mendekripsi data file menggunakan AES-256 dengan pemrosesan chunk untuk file besar.
    
    Parameter:
        encrypted_data (bytes): Data terenkripsi dengan IV di awal
    
    Return:
        bytes: Data asli yang sudah didekripsi
    
    Catatan:
        - Mengekstrak IV dari 16 byte pertama data terenkripsi
        - Untuk file besar (>10MB), memproses dalam chunk 1MB
        - Untuk file kecil, menggunakan pendekatan sederhana
        - Menggunakan PKCS7 unpadding untuk menghilangkan padding
        - Akan raise exception jika dekripsi gagal atau data tidak valid
    """
    try:
        if len(encrypted_data) < 16:
            raise ValueError("Encrypted data too short")
            
        # Extract IV from the beginning of the encrypted data
        iv = encrypted_data[:16]
        actual_encrypted_data = encrypted_data[16:]
        
        if len(actual_encrypted_data) == 0:
            raise ValueError("No encrypted data after IV")
        
        # For large files, use a more efficient approach with less memory overhead
        if len(actual_encrypted_data) > 10 * 1024 * 1024:  # 10MB
            # Process in chunks for large files
            chunk_size = 1024 * 1024  # 1MB chunks
            cipher = Cipher(algorithms.AES(ENCRYPTION_KEY), modes.CBC(iv), backend=backend)
            decryptor = cipher.decryptor()
            unpadder = padding.PKCS7(algorithms.AES.block_size).unpadder()
            
            # Process all but the last chunk
            result = bytearray()
            for i in range(0, len(actual_encrypted_data) - chunk_size, chunk_size):
                chunk = actual_encrypted_data[i:i+chunk_size]
                result.extend(decryptor.update(chunk))
            
            # Process the last chunk with unpadding
            last_chunk = actual_encrypted_data[-(len(actual_encrypted_data) % chunk_size or chunk_size):]
            decrypted_last_chunk = decryptor.update(last_chunk) + decryptor.finalize()
            result.extend(unpadder.update(decrypted_last_chunk) + unpadder.finalize())
            
            return bytes(result)
        else:
            # For smaller files, use the simpler approach
            cipher = Cipher(algorithms.AES(ENCRYPTION_KEY), modes.CBC(iv), backend=backend)
            decryptor = cipher.decryptor()
            padded_data = decryptor.update(actual_encrypted_data) + decryptor.finalize()
            
            unpadder = padding.PKCS7(algorithms.AES.block_size).unpadder()
            return unpadder.update(padded_data) + unpadder.finalize()
    except Exception as e:
        logger.error(f"AES decryption error: {str(e)}")
        raise

def decrypt_file_fernet(encrypted_data):
    """
    Mendekripsi data file menggunakan enkripsi Fernet.
    
    Parameter:
        encrypted_data (bytes): Data terenkripsi dengan Fernet
    
    Return:
        bytes: Data asli yang sudah didekripsi
    
    Catatan:
        - Fernet menangani IV dan padding secara otomatis
        - Digunakan sebagai fallback jika AES tidak tersedia
        - Akan raise exception jika dekripsi gagal
    """
    try:
        return cipher_suite.decrypt(encrypted_data)
    except Exception as e:
        logger.error(f"Fernet decryption error: {str(e)}")
        raise

def decrypt_file(encrypted_data):
    """
    Mendekripsi data file menggunakan metode terbaik yang tersedia.
    
    Parameter:
        encrypted_data (bytes): Data terenkripsi yang akan didekripsi
    
    Return:
        bytes: Data asli yang sudah didekripsi
    
    Catatan:
        - Prioritas: AES-256 (jika tersedia) > Fernet (fallback)
        - Mengukur waktu dekripsi untuk monitoring performa
        - Jika dekripsi gagal, mencoba menggunakan data as-is (mungkin tidak terenkripsi)
        - Akan raise ValueError jika semua metode dekripsi gagal
        - Log semua aktivitas dekripsi untuk debugging
    """
    start_time = time.time()
    try:
        # Try AES decryption first if it's available and the data looks like AES-encrypted
        if HAS_AES and len(encrypted_data) > 16:
            try:
                result = decrypt_file_aes(encrypted_data)
                decryption_time = time.time() - start_time
                logger.info(f"File decrypted successfully with AES in {decryption_time:.2f} seconds")
                return result
            except Exception:
                # If AES decryption fails, try Fernet
                if not HAS_AES:
                    result = decrypt_file_fernet(encrypted_data)
                    decryption_time = time.time() - start_time
                    logger.info(f"File decrypted successfully with Fernet in {decryption_time:.2f} seconds")
                    return result
        elif not HAS_AES:
            # If AES is not available, use Fernet
            result = decrypt_file_fernet(encrypted_data)
            decryption_time = time.time() - start_time
            logger.info(f"File decrypted successfully with Fernet in {decryption_time:.2f} seconds")
            return result
    except Exception as e:
        logger.error(f"Decryption error: {str(e)}")
        # If we can't decrypt, this might be unencrypted data (from fallback)
        try:
            # Check if it looks like valid file data
            if len(encrypted_data) > 0:
                logger.warning("Attempting to use data as-is (might be unencrypted)")
                return encrypted_data
        except:
            pass
        
        # If all fails, raise error
        raise ValueError("Failed to decrypt file data")

def session_expired(chat_id, file_path, db_id=None, lang='en'):
    """
    Menangani sesi yang telah berakhir (expired).
    
    Parameter:
        chat_id (int): ID chat Telegram
        file_path (str): Path file yang akan dihapus
        db_id (int, optional): ID database file untuk dihapus dari database
        lang (str): Kode bahasa untuk pesan ('en', 'id', 'ar', 'jv')
    
    Return:
        Tidak ada
    
    Catatan:
        - Menghapus file dari sistem jika ada
        - Menghapus record dari database jika db_id diberikan
        - Mengirim pesan expired ke pengguna dengan tombol restart
        - Menangani error dengan graceful untuk setiap operasi
    """
    try:
        # Delete the file if it exists
        if file_path and os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"Deleted expired file: {file_path}")
            except Exception as e:
                logger.error(f"Failed to delete expired file {file_path}: {str(e)}")
        
        # Delete from database if db_id is provided
        if db_id:
            try:
                conn = sqlite3.connect('files.db')
                conn.execute('DELETE FROM files WHERE id = ?', (db_id,))
                conn.commit()
                conn.close()
            except Exception as e:
                logger.error(f"Failed to delete expired file from database: {str(e)}")
        
        # Send expiration message
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['start_over'], callback_data="start_over"))
        bot.send_message(chat_id, LANG[lang]['session_expired'], reply_markup=markup)
        
    except Exception as e:
        logger.error(f"Error in session_expired: {str(e)}")

def update_countdown(chat_id, message_id, seconds_left, lang='en'):
    """
    Memperbarui pesan countdown timer.
    
    Parameter:
        chat_id (int): ID chat Telegram
        message_id (int): ID pesan yang akan diupdate
        seconds_left (int): Sisa waktu dalam detik
        lang (str): Kode bahasa untuk format pesan
    
    Return:
        Tidak ada
    
    Catatan:
        - Mengkonversi detik ke format menit:detik
        - Mengupdate pesan setiap detik untuk efek animasi
        - Menangani error jika update pesan gagal
    """
    try:
        # Format as minutes:seconds
        minutes = seconds_left // 60
        seconds = seconds_left % 60
        
        # Update every second for a more animated effect
        bot.edit_message_text(
            LANG[lang]['countdown'].format(minutes, seconds),
            chat_id,
            message_id
        )
    except Exception as e:
        logger.error(f"Error updating countdown: {str(e)}")

def start_session_timer(chat_id, file_path, db_id, lang='en'):
    """
    Memulai countdown timer untuk sesi pengguna.
    
    Parameter:
        chat_id (int): ID chat Telegram
        file_path (str): Path file yang akan dihapus saat expired
        db_id (int): ID database file
        lang (str): Kode bahasa untuk pesan timer
    
    Return:
        Tidak ada
    
    Catatan:
        - Membuat pesan countdown dengan format awal (2:00)
        - Menyimpan informasi sesi dalam active_sessions
        - Menggunakan threading.Timer untuk update setiap detik
        - Timer berjalan sebagai daemon thread
        - Otomatis memanggil session_expired() saat waktu habis
        - Menangani error dengan graceful
    """
    try:
        # Create countdown message with initial format (2:00)
        minutes = SESSION_TIMEOUT_SECONDS // 60
        seconds = SESSION_TIMEOUT_SECONDS % 60
        countdown_msg = bot.send_message(chat_id, LANG[lang]['countdown'].format(minutes, seconds))
        
        # Store session info
        active_sessions[chat_id] = {
            'file_path': file_path,
            'db_id': db_id,
            'countdown_msg_id': countdown_msg.message_id,
            'lang': lang,
            'start_time': time.time(),
            'timer': None
        }
        
        # Function to update countdown and check expiration
        def check_session():
            try:
                if chat_id not in active_sessions:
                    return
                
                elapsed = time.time() - active_sessions[chat_id]['start_time']
                remaining = max(0, SESSION_TIMEOUT_SECONDS - int(elapsed))
                
                # Always update the countdown for smooth animation
                
                if remaining <= 0:
                    # Session expired
                    session_info = active_sessions.pop(chat_id, None)
                    if session_info:
                        session_expired(chat_id, session_info['file_path'], session_info['db_id'], session_info['lang'])
                else:
                    # Update countdown and schedule next check - update every second for animation
                    update_countdown(chat_id, active_sessions[chat_id]['countdown_msg_id'], remaining, lang)
                    timer = threading.Timer(1.0, check_session)
                    timer.daemon = True
                    timer.start()
                    active_sessions[chat_id]['timer'] = timer
            except Exception as e:
                logger.error(f"Error in check_session: {str(e)}")
        
        # Start the timer
        timer = threading.Timer(1.0, check_session)
        timer.daemon = True
        timer.start()
        active_sessions[chat_id]['timer'] = timer
        
    except Exception as e:
        logger.error(f"Error starting session timer: {str(e)}")

def generate_secure_filename(original_name):
    """
    Menghasilkan nama file yang aman secara acak sambil mempertahankan ekstensi.
    
    Parameter:
        original_name (str): Nama file asli
    
    Return:
        str: Nama file aman dengan UUID dan ekstensi asli
    
    Catatan:
        - Menggunakan UUID4 untuk nama file yang unik dan aman
        - Mempertahankan ekstensi file asli jika ada
        - Mencegah path traversal dan karakter berbahaya
        - Format: {uuid}.{ekstensi} atau {uuid} jika tidak ada ekstensi
    """
    ext = original_name.split('.')[-1] if '.' in original_name else ''
    secure_name = str(uuid.uuid4())
    return f"{secure_name}.{ext}" if ext else secure_name

def secure_delete_file(file_path):
    """
    Menghapus file secara aman dengan menimpa data acak sebelum penghapusan.
    
    Parameter:
        file_path (str): Path file yang akan dihapus secara aman
    
    Return:
        Tidak ada
    
    Catatan:
        - Menimpa file dengan data acak sebelum dihapus untuk keamanan
        - File besar (>100MB) dihapus langsung untuk performa
        - Menggunakan secrets.token_bytes() untuk data acak yang aman
        - Melakukan fsync() untuk memastikan data tertulis ke disk
        - Fallback ke penghapusan biasa jika secure delete gagal
        - Last resort: truncate file ke 0 bytes jika semua gagal
    """
    if not os.path.exists(file_path):
        return
        
    try:
        # Get file size
        file_size = os.path.getsize(file_path)
        
        # Skip very large files for performance reasons
        if file_size > 100 * 1024 * 1024:  # 100MB
            os.remove(file_path)
            logger.info(f"Large file deleted (regular): {file_path}")
            return
            
        # Overwrite file with random data multiple times
        for _ in range(SECURE_DELETE_PASSES):
            with open(file_path, 'wb') as f:
                f.write(secrets.token_bytes(file_size))
                f.flush()
                os.fsync(f.fileno())
        
        # Finally delete the file
        os.remove(file_path)
        logger.info(f"Securely deleted: {file_path}")
    except Exception as e:
        logger.error(f"Secure deletion failed for {file_path}: {str(e)}")
        # Attempt regular deletion as fallback
        try:
            os.remove(file_path)
            logger.info(f"Fallback regular deletion for: {file_path}")
        except Exception as e2:
            logger.critical(f"Complete deletion failure for {file_path}: {str(e2)}")
            # Last resort - try to make file empty
            try:
                with open(file_path, 'w') as f:
                    pass  # Truncate file to 0 bytes
            except:
                pass

def cleanup_inactive_users():
    """
    Membersihkan sesi pengguna yang tidak aktif.
    
    Parameter:
        Tidak ada
    
    Return:
        Tidak ada
    
    Catatan:
        - Berjalan dalam loop tak terbatas sebagai background thread
        - Memeriksa aktivitas pengguna setiap menit
        - Menghapus pengguna yang tidak aktif lebih dari 5 menit
        - Membatalkan timer yang sedang berjalan untuk pengguna tidak aktif
        - Menangani error dengan graceful untuk menjaga stabilitas
    """
    while True:
        try:
            current_time = time.time()
            users_to_remove = []
            
            # Check each user's activity
            for user_id, data in user_activity.items():
                # If user has been inactive for more than 5 minutes, clean up
                if current_time - data['timestamp'] > 300:  # 5 minutes
                    users_to_remove.append(user_id)
            
            # Remove inactive users
            for user_id in users_to_remove:
                if user_id in user_activity:
                    try:
                        if user_activity[user_id]['timer']:
                            user_activity[user_id]['timer'].cancel()
                    except:
                        pass
                    user_activity.pop(user_id, None)
                    logger.info(f"Cleaned up inactive user: {user_id}")
        except Exception as e:
            logger.error(f"Error cleaning up inactive users: {str(e)}")
        
        # Check every minute
        time.sleep(60)

def cleanup_files():
    """
    Membersihkan file-file lama dan temporary secara berkala.
    
    Parameter:
        Tidak ada
    
    Return:
        Tidak ada
    
    Catatan:
        - Berjalan dalam loop tak terbatas sebagai background thread
        - Memeriksa dan membersihkan file setiap 3 menit
        - Menghapus file dari database yang lebih lama dari FILE_RETENTION_MINUTES
        - Membersihkan direktori temp dengan file lebih lama dari TEMP_FILE_RETENTION_MINUTES
        - Menghapus direktori kosong untuk menjaga kebersihan sistem
        - Menangani error dengan graceful untuk setiap operasi
    """
    while True:
        try:
            # Clean database files
            conn = sqlite3.connect('files.db', timeout=10.0)
            cutoff = datetime.now() - timedelta(minutes=FILE_RETENTION_MINUTES)
            cursor = conn.execute('SELECT file_path FROM files WHERE created_at < ?', (cutoff,))
            files_to_delete = cursor.fetchall()
            
            for (file_path,) in files_to_delete:
                if file_path and os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                        logger.info(f"Deleted old file: {file_path}")
                    except Exception as e:
                        logger.error(f"Failed to delete {file_path}: {str(e)}")
            
            conn.execute('DELETE FROM files WHERE created_at < ?', (cutoff,))
            conn.commit()
            conn.close()
            
            # Clean temp directory - more aggressive cleanup
            if os.path.exists("temp"):
                for filename in os.listdir("temp"):
                    file_path = os.path.join("temp", filename)
                    if os.path.isfile(file_path):
                        try:
                            file_age = datetime.now() - datetime.fromtimestamp(os.path.getctime(file_path))
                            if file_age > timedelta(minutes=TEMP_FILE_RETENTION_MINUTES):
                                os.remove(file_path)
                                logger.info(f"Deleted temp file: {file_path}")
                        except Exception as e:
                            logger.error(f"Failed to delete temp file {file_path}: {str(e)}")
            
            # Check for any empty directories and clean them
            for folder in ["files", "temp"]:
                if os.path.exists(folder) and os.path.isdir(folder):
                    try:
                        # Remove empty subdirectories
                        for root, dirs, files in os.walk(folder, topdown=False):
                            for dir_name in dirs:
                                dir_path = os.path.join(root, dir_name)
                                if not os.listdir(dir_path):  # If directory is empty
                                    os.rmdir(dir_path)
                                    logger.info(f"Removed empty directory: {dir_path}")
                    except Exception as e:
                        logger.error(f"Failed to clean directories in {folder}: {str(e)}")
                        
        except Exception as e:
            logger.error(f"Cleanup error: {str(e)}")
        time.sleep(180)  # Check every 3 minutes instead of 5

# Start background cleanup threads
threading.Thread(target=cleanup_files, daemon=True).start()
threading.Thread(target=cleanup_inactive_users, daemon=True).start()

# Start PDF merge session cleanup
def cleanup_merge_sessions():
    """
    Membersihkan sesi PDF merge yang sudah expired.
    
    Parameter:
        Tidak ada
    
    Return:
        Tidak ada
    
    Catatan:
        - Berjalan dalam loop tak terbatas sebagai background thread
        - Memeriksa sesi merge yang expired setiap menit
        - Memanggil cleanup_expired_merge_sessions() untuk pembersihan aktual
        - Menangani error dengan graceful untuk menjaga stabilitas
    """
    while True:
        try:
            cleanup_expired_merge_sessions()
        except Exception as e:
            logger.error(f"Error in merge session cleanup: {str(e)}")
        time.sleep(60)  # Check every minute

threading.Thread(target=cleanup_merge_sessions, daemon=True).start()

@bot.message_handler(func=lambda message: message.content_type == 'text' and not message.text.startswith('/'))
def handle_first_message(message):
    """
    Menangani pesan teks pertama sebagai interaksi awal.
    
    Parameter:
        message: Objek pesan Telegram
    
    Return:
        Tidak ada
    
    Catatan:
        - Handler untuk semua pesan teks yang bukan command (/)
        - Mengirim pesan welcome pertama dengan tombol Start
        - Memulai timer inactivity untuk pengguna
        - Mendeteksi bahasa pengguna dari language_code Telegram
        - Menangani error dengan graceful
    """
    try:
        user_id = message.from_user.id
        lang = get_user_lang(message.from_user.language_code)
        
        # Send first welcome message with Start button
        send_first_welcome(message.chat.id, lang)
        
        # Start inactivity timer
        start_inactivity_timer(user_id, message.chat.id, lang)
        
    except Exception as e:
        logger.error(f"Error in handle_first_message: {str(e)}")

def send_first_welcome(chat_id, lang='en'):
    """
    Mengirim pesan welcome pertama dengan tombol Start.
    
    Parameter:
        chat_id (int): ID chat Telegram
        lang (str): Kode bahasa untuk pesan ('en', 'id', 'ar', 'jv')
    
    Return:
        Tidak ada
    
    Catatan:
        - Membuat inline keyboard dengan tombol Start
        - Menggunakan teks welcome dari dictionary LANG
        - Menangani error jika pengiriman pesan gagal
    """
    try:
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['start_button'], callback_data="start_bot"))
        bot.send_message(chat_id, LANG[lang]['first_welcome'], reply_markup=markup)
    except Exception as e:
        logger.error(f"Error sending first welcome: {str(e)}")

def send_error_with_restart(chat_id, error_message, lang='en'):
    """
    Mengirim pesan error dengan tombol Back to Start.
    
    Parameter:
        chat_id (int): ID chat Telegram
        error_message (str): Pesan error yang akan ditampilkan
        lang (str): Kode bahasa untuk tombol dan pesan tambahan
    
    Return:
        Tidak ada
    
    Catatan:
        - Membuat inline keyboard dengan tombol "Back to Start"
        - Menggabungkan pesan error dengan instruksi try_again
        - Memberikan cara mudah bagi pengguna untuk restart bot
        - Menangani error jika pengiriman pesan gagal
    """
    try:
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_start'], callback_data="start_bot"))
        
        full_message = f"{error_message}\n\n{LANG[lang]['try_again']}"
        bot.send_message(chat_id, full_message, reply_markup=markup)
    except Exception as e:
        logger.error(f"Error sending error message: {str(e)}")

def cleanup_failed_file(file_path):
    """
    Membersihkan file yang gagal atau rusak.
    
    Parameter:
        file_path (str): Path file yang akan dibersihkan
    
    Return:
        Tidak ada
    
    Catatan:
        - Memeriksa apakah file ada sebelum mencoba menghapus
        - Menghapus file yang gagal diproses atau rusak
        - Menangani error jika penghapusan gagal
        - Digunakan untuk cleanup setelah error processing
    """
    try:
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
            logger.info(f"Cleaned up failed file: {file_path}")
    except Exception as e:
        logger.error(f"Failed to clean up file {file_path}: {str(e)}")

def start_inactivity_timer(user_id, chat_id, lang='en'):
    """
    Memulai timer inactivity untuk pengguna.
    
    Parameter:
        user_id (int): ID pengguna Telegram
        chat_id (int): ID chat Telegram
        lang (str): Kode bahasa untuk pesan
    
    Return:
        Tidak ada
    
    Catatan:
        - Menyimpan informasi aktivitas pengguna dalam user_activity
        - Memeriksa inactivity setiap 30 detik
        - Mengirim reminder setelah 2 menit tidak aktif
        - Menutup sesi setelah 3 menit tidak aktif
        - Menggunakan threading.Timer untuk penjadwalan
        - Timer berjalan sebagai daemon thread
    """
    # Update user's last activity time
    user_activity[user_id] = {
        'timestamp': time.time(),
        'chat_id': chat_id,
        'lang': lang,
        'reminder_sent': False,
        'timer': None
    }
    
    # Function to check inactivity
    def check_inactivity():
        if user_id not in user_activity:
            return
            
        current_time = time.time()
        last_activity = user_activity[user_id]['timestamp']
        elapsed = current_time - last_activity
        
        # If 2 minutes passed without activity and reminder not sent yet
        if elapsed > 120 and not user_activity[user_id]['reminder_sent']:
            # Send reminder
            try:
                bot.send_message(chat_id, LANG[lang]['inactivity_reminder'])
                user_activity[user_id]['reminder_sent'] = True
                
                # Schedule final check after 1 more minute
                timer = threading.Timer(60.0, check_inactivity)
                timer.daemon = True
                timer.start()
                user_activity[user_id]['timer'] = timer
            except:
                pass
        # If 3 minutes total passed (reminder sent 1 minute ago)
        elif elapsed > 180 and user_activity[user_id]['reminder_sent']:
            # Send session closed message
            try:
                markup = types.InlineKeyboardMarkup()
                markup.add(types.InlineKeyboardButton(LANG[lang]['start_button'], callback_data="start_bot"))
                bot.send_message(chat_id, LANG[lang]['inactivity_close'], reply_markup=markup)
                
                # Clean up user activity
                user_activity.pop(user_id, None)
            except:
                pass
        else:
            # Schedule next check
            timer = threading.Timer(30.0, check_inactivity)  # Check every 30 seconds
            timer.daemon = True
            timer.start()
            user_activity[user_id]['timer'] = timer
    
    # Start the inactivity timer
    timer = threading.Timer(30.0, check_inactivity)  # First check after 30 seconds
    timer.daemon = True
    timer.start()
    user_activity[user_id]['timer'] = timer

def update_user_activity(user_id):
    """
    Memperbarui timestamp aktivitas terakhir pengguna.
    
    Parameter:
        user_id (int): ID pengguna Telegram
    
    Return:
        Tidak ada
    
    Catatan:
        - Memperbarui timestamp ke waktu saat ini
        - Reset flag reminder_sent ke False
        - Digunakan setiap kali pengguna melakukan aktivitas
        - Mencegah pengiriman reminder yang tidak perlu
    """
    if user_id in user_activity:
        user_activity[user_id]['timestamp'] = time.time()
        user_activity[user_id]['reminder_sent'] = False

def security_check_user(user_id):
    """
    Memeriksa apakah pengguna terkena rate limit atau diblokir.
    
    Parameter:
        user_id (int): ID pengguna Telegram
    
    Return:
        bool: True jika pengguna diizinkan, False jika diblokir
    
    Catatan:
        - Memeriksa apakah pengguna sudah diblokir sebelumnya
        - Membersihkan request lama di luar time window
        - Menghitung jumlah request dalam window waktu
        - Memblokir pengguna jika melebihi RATE_LIMIT_REQUESTS
        - Mencatat request baru untuk tracking
        - Menggunakan sliding window untuk rate limiting
    """
    if user_id in blocked_users:
        return False
    
    current_time = time.time()
    if user_id not in user_request_counts:
        user_request_counts[user_id] = []
    
    # Clean old requests outside time window
    user_request_counts[user_id] = [req_time for req_time in user_request_counts[user_id] 
                                   if current_time - req_time < RATE_LIMIT_WINDOW]
    
    # Check rate limit
    if len(user_request_counts[user_id]) >= RATE_LIMIT_REQUESTS:
        blocked_users.add(user_id)
        logger.warning(f"User {user_id} blocked for rate limiting")
        return False
    
    # Add current request
    user_request_counts[user_id].append(current_time)
    return True

def validate_file_security(message):
    """
    Validasi keamanan file yang ditingkatkan.
    
    Parameter:
        message: Objek pesan Telegram yang berisi file
    
    Return:
        tuple: (is_valid, error_message)
               - is_valid (bool): True jika file aman, False jika tidak
               - error_message (str): Pesan error jika validasi gagal
    
    Catatan:
        - Memeriksa ukuran file (maksimal MAX_FILE_SIZE)
        - Memvalidasi ekstensi file terhadap ALLOWED_FILE_TYPES
        - Memeriksa MIME type untuk keamanan tambahan
        - Mendukung dokumen, gambar, audio, dan video
        - Mengembalikan pesan error yang spesifik untuk debugging
    """
    if message.content_type == 'document':
        file_size = message.document.file_size
        file_name = message.document.file_name
        mime_type = message.document.mime_type
        
        # Check file size
        if file_size > MAX_FILE_SIZE:
            return False, "File too large"
        
        # Check file extension
        if file_name:
            ext = file_name.lower().split('.')[-1] if '.' in file_name else ''
            if ext not in ALLOWED_FILE_TYPES:
                return False, "File type not allowed"
        
        # Check MIME type for PDFs and DOCX
        if mime_type and mime_type not in ['application/pdf', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document', 'application/msword']:
            if not mime_type.startswith('image/') and not mime_type.startswith('audio/') and not mime_type.startswith('video/'):
                return False, "Invalid MIME type"
    
    return True, "OK"

def validate_file_for_service(message, service, lang):
    """
    Memvalidasi apakah file yang diupload sesuai dengan layanan yang dipilih.
    
    Parameter:
        message: Objek pesan Telegram yang berisi file
        service (str): Nama layanan yang dipilih pengguna
        lang (str): Kode bahasa untuk pesan error
    
    Return:
        bool: True jika file sesuai layanan, False jika tidak
    
    Catatan:
        - Melakukan validasi keamanan terlebih dahulu
        - Menentukan jenis file berdasarkan content_type dan nama file
        - Memeriksa kompatibilitas file dengan layanan:
          * PDF services: hanya menerima file PDF
          * Image services: hanya menerima file gambar
          * Media services: hanya menerima audio/video
          * Document/ZIP: menerima semua file yang didukung
        - Mengirim pesan error dengan tombol kembali jika tidak sesuai
    """
    # Security validation first
    is_secure, error_msg = validate_file_security(message)
    if not is_secure:
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
        bot.reply_to(message, f"❌ Security check failed: {error_msg}", reply_markup=markup)
        return False
    
    file_type = None
    ext = None
    
    if message.content_type == 'document':
        file_type, ext = get_file_type(message.document.file_name)
    elif message.content_type == 'photo':
        file_type, ext = 'image', 'jpg'
    elif message.content_type == 'video':
        file_type, ext = 'video', 'mp4'
    elif message.content_type == 'audio':
        file_type, ext = 'audio', 'mp3'
    
    # Check service compatibility
    if service in ['pdf_merge', 'pdf_compress', 'pdf_convert', 'pdf_split'] and (file_type != 'document' or ext != 'pdf'):
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
        bot.reply_to(message, LANG[lang]['wrong_file_type'], reply_markup=markup)
        return False
    elif service in ['image', 'compress_image'] and file_type != 'image':
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
        bot.reply_to(message, LANG[lang]['wrong_file_type'], reply_markup=markup)
        return False
    elif service == 'media' and file_type not in ['audio', 'video']:
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
        bot.reply_to(message, LANG[lang]['wrong_file_type'], reply_markup=markup)
        return False
    elif service in ['document', 'compress_zip'] and file_type == 'unsupported':
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
        bot.reply_to(message, LANG[lang]['wrong_file_type'], reply_markup=markup)
        return False
    
    return True

def create_pdf_merge_session(user_id, chat_id, lang='en'):
    """
    Membuat sesi PDF merge baru untuk pengguna.
    
    Parameter:
        user_id (int): ID pengguna Telegram
        chat_id (int): ID chat Telegram
        lang (str): Kode bahasa untuk pesan
    
    Return:
        Tidak ada
    
    Catatan:
        - Membuat entry baru dalam pdf_merge_sessions
        - Menyimpan informasi: chat_id, daftar PDF, bahasa, waktu dibuat
        - Mengatur flag awaiting_files untuk batch collection
        - Memulai timer 5 detik untuk mengumpulkan multiple files
        - Timer otomatis menampilkan konfirmasi urutan setelah batch selesai
    """
    pdf_merge_sessions[user_id] = {
        'chat_id': chat_id,
        'pdfs': [],
        'lang': lang,
        'created_at': time.time(),
        'awaiting_files': True,
        'batch_timer': None
    }
    
    # Start batch collection timer (5 seconds to collect multiple files)
    def end_batch_collection():
        if user_id in pdf_merge_sessions and pdf_merge_sessions[user_id]['awaiting_files']:
            pdf_merge_sessions[user_id]['awaiting_files'] = False
            show_pdf_order_confirmation(user_id)
    
    timer = threading.Timer(5.0, end_batch_collection)
    timer.daemon = True
    timer.start()
    pdf_merge_sessions[user_id]['batch_timer'] = timer

def add_pdf_to_merge_session(user_id, pdf_id):
    """
    Menambahkan PDF ke sesi merge yang sudah ada.
    
    Parameter:
        user_id (int): ID pengguna Telegram
        pdf_id (int): ID database dari file PDF
    
    Return:
        bool: True jika berhasil ditambahkan, False jika gagal
    
    Catatan:
        - Memeriksa apakah sesi merge ada untuk pengguna
        - Membatasi maksimal 10 PDF per sesi
        - Menambahkan PDF ID ke daftar pdfs dalam sesi
        - Reset batch timer jika masih dalam mode collecting
        - Timer baru 3 detik setelah file terakhir ditambahkan
        - Otomatis menampilkan konfirmasi setelah batch selesai
    """
    if user_id in pdf_merge_sessions:
        if len(pdf_merge_sessions[user_id]['pdfs']) < 10:  # Limit to 10 PDFs
            pdf_merge_sessions[user_id]['pdfs'].append(pdf_id)
            # Reset batch timer if still collecting
            if pdf_merge_sessions[user_id]['awaiting_files']:
                if pdf_merge_sessions[user_id]['batch_timer']:
                    pdf_merge_sessions[user_id]['batch_timer'].cancel()
                
                def end_batch_collection():
                    if user_id in pdf_merge_sessions and pdf_merge_sessions[user_id]['awaiting_files']:
                        pdf_merge_sessions[user_id]['awaiting_files'] = False
                        show_pdf_order_confirmation(user_id)
                
                timer = threading.Timer(3.0, end_batch_collection)  # 3 seconds after last file
                timer.daemon = True
                timer.start()
                pdf_merge_sessions[user_id]['batch_timer'] = timer
            return True
    return False

def get_pdf_merge_session(user_id):
    """
    Mengambil sesi PDF merge untuk pengguna.
    
    Parameter:
        user_id (int): ID pengguna Telegram
    
    Return:
        dict atau None: Data sesi merge jika ada, None jika tidak ada
    
    Catatan:
        - Menggunakan dict.get() untuk menghindari KeyError
        - Mengembalikan None jika pengguna tidak memiliki sesi aktif
        - Data sesi berisi: chat_id, pdfs, lang, created_at, awaiting_files, batch_timer
    """
    return pdf_merge_sessions.get(user_id)

def clear_pdf_merge_session(user_id):
    """
    Membersihkan sesi PDF merge untuk pengguna.
    
    Parameter:
        user_id (int): ID pengguna Telegram
    
    Return:
        Tidak ada
    
    Catatan:
        - Membatalkan batch timer jika masih aktif
        - Menghapus semua file PDF dari sistem dan database
        - Menggunakan cleanup_failed_file() untuk penghapusan aman
        - Menghapus entry sesi dari pdf_merge_sessions
        - Menangani error untuk setiap operasi cleanup
        - Memastikan tidak ada file yang tertinggal
    """
    if user_id in pdf_merge_sessions:
        # Clean up PDF files
        session = pdf_merge_sessions[user_id]
        
        # Cancel batch timer if active
        if session.get('batch_timer'):
            try:
                session['batch_timer'].cancel()
            except:
                pass
        
        for pdf_id in session['pdfs']:
            try:
                conn = sqlite3.connect('files.db')
                cursor = conn.execute('SELECT file_path FROM files WHERE id = ?', (pdf_id,))
                result = cursor.fetchone()
                if result:
                    cleanup_failed_file(result[0])
                    conn.execute('DELETE FROM files WHERE id = ?', (pdf_id,))
                    conn.commit()
                conn.close()
            except Exception as e:
                logger.error(f"Error cleaning up PDF {pdf_id}: {str(e)}")
        
        pdf_merge_sessions.pop(user_id, None)

def generate_pdf_list_text(user_id, lang='en'):
    """
    Menghasilkan teks yang menampilkan daftar PDF saat ini.
    
    Parameter:
        user_id (int): ID pengguna Telegram
        lang (str): Kode bahasa (tidak digunakan saat ini)
    
    Return:
        str: Teks berisi daftar PDF dengan nomor urut
    
    Catatan:
        - Mengambil sesi merge untuk pengguna
        - Mengembalikan string kosong jika tidak ada sesi
        - Mengquery database untuk mendapatkan nama file
        - Format: "1. nama_file.pdf\n2. nama_file2.pdf"
        - Menangani koneksi database dengan aman
    """
    session = get_pdf_merge_session(user_id)
    if not session:
        return ""
    
    text_lines = []
    conn = sqlite3.connect('files.db')
    
    for i, pdf_id in enumerate(session['pdfs'], 1):
        cursor = conn.execute('SELECT file_name FROM files WHERE id = ?', (pdf_id,))
        result = cursor.fetchone()
        if result:
            filename = result[0]
            text_lines.append(f"{i}. {filename}")
    
    conn.close()
    return "\n".join(text_lines)

def create_pdf_reorder_markup(user_id, lang='en'):
    """
    Membuat inline keyboard untuk pengurutan ulang PDF.
    
    Parameter:
        user_id (int): ID pengguna Telegram
        lang (str): Kode bahasa untuk tombol
    
    Return:
        InlineKeyboardMarkup: Keyboard dengan tombol pengurutan
    
    Catatan:
        - Mengembalikan tombol cancel saja jika PDF kurang dari 2
        - Membuat tombol "Merge Now" di atas
        - Untuk setiap PDF: tombol naik, nama file, tombol turun, tombol hapus
        - Menggunakan placeholder "➖" untuk tombol yang tidak aktif
        - Memotong nama file jika lebih dari 15 karakter
        - Menambahkan tombol kembali dan cancel di bawah
    """
    session = get_pdf_merge_session(user_id)
    if not session or len(session['pdfs']) < 2:
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['cancel_merge'], callback_data="cancel_pdf_merge"))
        return markup
    
    markup = types.InlineKeyboardMarkup()
    
    # Add merge button
    markup.add(types.InlineKeyboardButton('🔗 Merge Now', callback_data="execute_pdf_merge"))
    
    # Add reorder buttons for each PDF
    conn = sqlite3.connect('files.db')
    for i, pdf_id in enumerate(session['pdfs']):
        cursor = conn.execute('SELECT file_name FROM files WHERE id = ?', (pdf_id,))
        result = cursor.fetchone()
        if result:
            filename = result[0][:15] + "..." if len(result[0]) > 15 else result[0]
            row = []
            
            # Move up button (not for first item)
            if i > 0:
                row.append(types.InlineKeyboardButton("⬆️", callback_data=f"move_pdf_up_{i}"))
            else:
                row.append(types.InlineKeyboardButton("➖", callback_data="noop"))  # Placeholder
            
            # PDF name with position
            row.append(types.InlineKeyboardButton(f"{i+1}. {filename}", callback_data="noop"))
            
            # Move down button (not for last item)
            if i < len(session['pdfs']) - 1:
                row.append(types.InlineKeyboardButton("⬇️", callback_data=f"move_pdf_down_{i}"))
            else:
                row.append(types.InlineKeyboardButton("➖", callback_data="noop"))  # Placeholder
            
            # Remove button
            row.append(types.InlineKeyboardButton("❌", callback_data=f"remove_pdf_{i}"))
            
            markup.row(*row)
    
    conn.close()
    
    # Add back to confirmation and cancel buttons
    markup.add(types.InlineKeyboardButton('↩️ Back to Confirmation', callback_data="back_to_confirmation"))
    markup.add(types.InlineKeyboardButton(LANG[lang]['cancel_merge'], callback_data="cancel_pdf_merge"))
    
    return markup

def merge_pdfs(user_id, lang='en'):
    """
    Menggabungkan PDF-PDF dalam sesi.
    
    Parameter:
        user_id (int): ID pengguna Telegram
        lang (str): Kode bahasa untuk pesan error
    
    Return:
        tuple: (pdf_data, error_message)
               - pdf_data (bytes): Data PDF yang sudah digabung, None jika gagal
               - error_message (str): Pesan error jika gagal, None jika berhasil
    
    Catatan:
        - Memeriksa apakah ada minimal 2 PDF dalam sesi
        - Menggunakan PdfMerger untuk menggabungkan PDF
        - Membaca dan mendekripsi setiap PDF dari database
        - Membuat temporary BytesIO untuk setiap PDF
        - Mengembalikan data PDF gabungan dalam bytes
        - Membersihkan semua resource dengan aman di finally block
    """
    session = get_pdf_merge_session(user_id)
    if not session or len(session['pdfs']) < 2:
        return None, LANG[lang]['pdf_merge_min_files']
    
    merger = None
    conn = None
    temp_pdfs = []
    
    try:
        merger = PdfMerger()
        conn = sqlite3.connect('files.db')
        
        # Add each PDF to merger
        for pdf_id in session['pdfs']:
            cursor = conn.execute('SELECT file_path FROM files WHERE id = ?', (pdf_id,))
            result = cursor.fetchone()
            if result:
                file_path = result[0]
                # Read and decrypt PDF
                with open(file_path, 'rb') as f:
                    encrypted_data = f.read()
                pdf_data = decrypt_file(encrypted_data)
                
                # Create temporary file for merger
                temp_pdf = BytesIO(pdf_data)
                temp_pdfs.append(temp_pdf)
                merger.append(temp_pdf)
        
        # Create output
        output = BytesIO()
        merger.write(output)
        output.seek(0)
        
        return output.getvalue(), None
        
    except Exception as e:
        logger.error(f"PDF merge error: {str(e)}")
        return None, LANG[lang]['pdf_merge_failed']
    finally:
        # Clean up resources
        if merger:
            try:
                merger.close()
            except:
                pass
        if conn:
            try:
                conn.close()
            except:
                pass
        for temp_pdf in temp_pdfs:
            try:
                temp_pdf.close()
            except:
                pass

def show_pdf_order_confirmation(user_id):
    """
    Menampilkan konfirmasi urutan PDF setelah batch collection.
    
    Parameter:
        user_id (int): ID pengguna Telegram
    
    Return:
        Tidak ada
    
    Catatan:
        - Mengambil sesi merge dan memeriksa apakah ada PDF
        - Menampilkan pesan error jika PDF kurang dari 2
        - Menghasilkan daftar PDF dengan urutan default
        - Menampilkan konfirmasi dengan 3 opsi:
          * Merge Now (Keep Order): langsung gabung dengan urutan saat ini
          * Change Order: buka interface pengurutan ulang
          * Cancel Merge: batalkan proses merge
        - Menggunakan Markdown untuk formatting pesan
    """
    session = get_pdf_merge_session(user_id)
    if not session or len(session['pdfs']) == 0:
        return
    
    lang = session['lang']
    chat_id = session['chat_id']
    
    if len(session['pdfs']) < 2:
        # Need at least 2 PDFs
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton(LANG[lang]['cancel_merge'], callback_data="cancel_pdf_merge"))
        bot.send_message(chat_id, LANG[lang]['pdf_merge_min_files'], reply_markup=markup)
        return
    
    # Show file order confirmation
    pdf_list = generate_pdf_list_text(user_id, lang)
    confirmation_text = f"📄 **Files received ({len(session['pdfs'])} PDFs):**\n\n{pdf_list}\n\n✅ **Default merge order shown above.**\n\n❓ Do you want to change the order before merging?"
    
    markup = types.InlineKeyboardMarkup()
    markup.add(types.InlineKeyboardButton('🔗 Merge Now (Keep Order)', callback_data="execute_pdf_merge"))
    markup.add(types.InlineKeyboardButton('✏️ Change Order', callback_data="show_reorder_options"))
    markup.add(types.InlineKeyboardButton(LANG[lang]['cancel_merge'], callback_data="cancel_pdf_merge"))
    
    bot.send_message(chat_id, confirmation_text, parse_mode='Markdown', reply_markup=markup)

def cleanup_expired_merge_sessions():
    """
    Membersihkan sesi PDF merge yang sudah expired.
    
    Parameter:
        Tidak ada
    
    Return:
        Tidak ada
    
    Catatan:
        - Memeriksa semua sesi merge yang aktif
        - Menandai sesi yang sudah lebih dari 5 menit (300 detik)
        - Mengirim pesan timeout ke pengguna yang sesinya expired
        - Memanggil clear_pdf_merge_session() untuk pembersihan
        - Menangani error untuk setiap operasi cleanup
        - Dipanggil secara berkala oleh background thread
    """
    current_time = time.time()
    expired_users = []
    
    for user_id, session in pdf_merge_sessions.items():
        if current_time - session['created_at'] > 300:  # 5 minutes
            expired_users.append(user_id)
    
    for user_id in expired_users:
        try:
            session = pdf_merge_sessions[user_id]
            bot.send_message(session['chat_id'], LANG[session['lang']]['pdf_merge_timeout'])
            clear_pdf_merge_session(user_id)
        except Exception as e:
            logger.error(f"Error cleaning expired merge session: {str(e)}")

@bot.message_handler(commands=['start'])
def start_message(message):
    """
    Handler untuk command /start - menampilkan menu utama bot.
    
    Parameter:
        message: Objek pesan Telegram yang berisi command /start
    
    Return:
        Tidak ada
    
    Catatan:
        - Melakukan security check untuk rate limiting
        - Membatalkan sesi aktif yang ada untuk pengguna
        - Membersihkan pilihan layanan sebelumnya
        - Menampilkan menu welcome dengan 4 layanan utama
        - Memulai timer inactivity untuk pengguna
        - Mendeteksi bahasa pengguna dan menyesuaikan pesan
        - Log aktivitas pengguna untuk monitoring
        - Mencegah duplikasi pesan dengan cleanup session
    """
    try:
        user_id = message.from_user.id
        username = message.from_user.username or 'Unknown'
        
        # Security check
        if not security_check_user(user_id):
            bot.send_message(message.chat.id, "❌ Access denied. Too many requests.")
            return
        
        logger.info(f"User {user_id} ({username}) started the bot")
        
        # Cancel any existing session for this user to prevent duplication
        if user_id in active_sessions:
            try:
                active_sessions[user_id]['timer'].cancel()
                # Try to delete countdown message if exists
                try:
                    bot.delete_message(message.chat.id, active_sessions[user_id]['countdown_msg_id'])
                except:
                    pass
            except:
                pass
            active_sessions.pop(user_id, None)
        
        # Clear any service selection and PDF merge sessions
        user_services.pop(user_id, None)
        if user_id in pdf_merge_sessions:
            clear_pdf_merge_session(user_id)
        
        # Clear user activity to prevent duplicate timers
        if user_id in user_activity:
            try:
                if user_activity[user_id]['timer']:
                    user_activity[user_id]['timer'].cancel()
            except:
                pass
            user_activity.pop(user_id, None)
        
        # Send welcome message with comprehensive service menu
        lang = get_user_lang(message.from_user.language_code)
        markup = types.InlineKeyboardMarkup()
        markup.add(types.InlineKeyboardButton('📄 PDF Tools', callback_data="service_pdf"))
        markup.add(types.InlineKeyboardButton('📸 Image Tools', callback_data="service_image"))
        markup.add(types.InlineKeyboardButton('🎵 Media Tools', callback_data="service_media"))
        markup.add(types.InlineKeyboardButton('🗜️ Compression', callback_data="service_compress"))
        
        bot.send_message(message.chat.id, LANG[lang]['welcome'], parse_mode='Markdown', reply_markup=markup)
        
        # Update user activity and start inactivity timer
        update_user_activity(user_id)
        start_inactivity_timer(user_id, message.chat.id, lang)
        
        logger.info(f"Welcome message sent to user {user_id} in {lang} language")
    except Exception as e:
        logger.error(f"Error in start_message: {str(e)}", exc_info=True)

def async_encrypt_file(file_data, file_path):
    """
    Mengenkripsi dan menyimpan file secara asinkron dengan optimasi memori.
    
    Parameter:
        file_data (bytes): Data file yang akan dienkripsi
        file_path (str): Path tempat file akan disimpan
    
    Return:
        float: Waktu yang dibutuhkan untuk enkripsi dalam detik
    
    Catatan:
        - Untuk file besar (>50MB), enkripsi langsung ke disk untuk menghemat memori
        - Menggunakan chunk 5MB untuk file sangat besar
        - Untuk file kecil, menggunakan enkripsi in-memory
        - Membuat file temporary saat enkripsi direct-to-disk
        - Mengukur waktu enkripsi untuk monitoring performa
        - Akan raise exception jika enkripsi gagal
    """
    try:
        # Start timing for performance measurement
        start_time = time.time()
        
        # For very large files, encrypt directly to disk to reduce memory usage
        if len(file_data) > 50 * 1024 * 1024:  # 50MB
            # Use direct-to-disk encryption for very large files
            if HAS_AES:
                # Create a temporary file for encryption
                temp_path = f"{file_path}.tmp"
                
                # Encrypt in chunks and write directly to disk
                chunk_size = 5 * 1024 * 1024  # 5MB chunks
                padder = padding.PKCS7(algorithms.AES.block_size).padder()
                cipher = Cipher(algorithms.AES(ENCRYPTION_KEY), modes.CBC(ENCRYPTION_IV), backend=backend)
                encryptor = cipher.encryptor()
                
                # Write IV at the beginning of the file
                with open(temp_path, 'wb') as f:
                    f.write(ENCRYPTION_IV)
                    
                    # Process all chunks except the last one
                    for i in range(0, len(file_data) - chunk_size, chunk_size):
                        chunk = file_data[i:i+chunk_size]
                        encrypted_chunk = encryptor.update(chunk)
                        f.write(encrypted_chunk)
                    
                    # Process the last chunk with padding
                    last_chunk = file_data[-(len(file_data) % chunk_size or chunk_size):]
                    padded_last_chunk = padder.update(last_chunk) + padder.finalize()
                    encrypted_last_chunk = encryptor.update(padded_last_chunk) + encryptor.finalize()
                    f.write(encrypted_last_chunk)
                
                # Rename the temp file to the final file
                os.replace(temp_path, file_path)
            else:
                # Use Fernet for large files (less efficient but simpler)
                encrypted_data = cipher_suite.encrypt(file_data)
                with open(file_path, 'wb') as f:
                    f.write(encrypted_data)
        else:
            # For smaller files, use the in-memory encryption
            encrypted_data, _ = encrypt_file(file_data)
            with open(file_path, 'wb') as f:
                f.write(encrypted_data)
        
        # Calculate total encryption time
        encryption_time = time.time() - start_time
        logger.info(f"File encrypted and saved in {encryption_time:.2f} seconds")
        return encryption_time
    except Exception as e:
        logger.error(f"Async encryption error: {str(e)}")
        raise

@bot.message_handler(content_types=['document', 'photo', 'video', 'audio'])
def handle_file(message):
    """
    Handler utama untuk memproses file yang dikirim pengguna.
    
    Parameter:
        message: Objek pesan Telegram yang berisi file (document, photo, video, audio)
    
    Return:
        Tidak ada
    
    Catatan:
        - Melakukan security check dan rate limiting
        - Memvalidasi apakah pengguna sudah memilih layanan
        - Menangani mode PDF merge dengan batch collection
        - Mengenkripsi file secara asinkron untuk keamanan
        - Menyimpan file ke database dengan nama aman
        - Menampilkan menu kontekstual berdasarkan jenis file dan layanan
        - Memulai session timer untuk keamanan
        - Menangani berbagai jenis error dengan graceful
    """
    try:
        user_id = message.from_user.id
        lang = get_user_lang(message.from_user.language_code)
        
        # Security check
        if not security_check_user(user_id):
            bot.reply_to(message, "❌ Access denied. Too many requests.")
            return
        
        # Check if user has selected a service
        if user_id not in user_services:
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
            bot.reply_to(message, "Please choose a service first by typing /start", reply_markup=markup)
            return
        
        service = user_services[user_id]
        
        # Check if user is in PDF merge mode
        merge_session = get_pdf_merge_session(user_id)
        if merge_session:
            # Handle PDF merge mode
            if message.content_type == 'document':
                file_type, ext = get_file_type(message.document.file_name)
                if ext == 'pdf':
                    # Add PDF to merge session
                    if len(merge_session['pdfs']) >= 10:
                        bot.reply_to(message, LANG[lang]['pdf_merge_limit'])
                        return
                    
                    # Process the PDF file
                    file_info = bot.get_file(message.document.file_id)
                    original_name = message.document.file_name
                    secure_filename = generate_secure_filename(original_name)
                    downloaded_file = bot.download_file(file_info.file_path)
                    file_path = f"files/{secure_filename}"
                    
                    # Validate PDF
                    try:
                        if HAS_PDF_MERGER:
                            temp_pdf = BytesIO(downloaded_file)
                            test_merger = PdfMerger()
                            test_merger.append(temp_pdf)
                            test_merger.close()
                            temp_pdf.close()
                    except Exception:
                        bot.reply_to(message, LANG[lang]['pdf_file_corrupted'])
                        return
                    
                    # Encrypt and store
                    encrypted_data, _ = encrypt_file(downloaded_file)
                    with open(file_path, 'wb') as f:
                        f.write(encrypted_data)
                    
                    # Store in database
                    try:
                        conn = sqlite3.connect('files.db', timeout=10.0)
                        cursor = conn.execute('INSERT INTO files (user_id, file_id, file_name, file_path, created_at) VALUES (?, ?, ?, ?, ?)',
                                    (user_id, file_info.file_id, original_name, file_path, datetime.now()))
                        db_id = cursor.lastrowid
                        conn.commit()
                        conn.close()
                    except Exception as db_error:
                        logger.error(f"Database insert failed for PDF merge: {str(db_error)}")
                        cleanup_failed_file(file_path)
                        bot.reply_to(message, LANG[lang]['error_upload'])
                        return
                    
                    # Add to merge session
                    add_pdf_to_merge_session(user_id, db_id)
                    
                    # Show brief confirmation (batch collection in progress)
                    if merge_session['awaiting_files']:
                        bot.reply_to(message, f"✅ **PDF {len(merge_session['pdfs'])}:** {original_name}\n\n⏱️ Send more PDFs or wait for batch collection to finish...", parse_mode='Markdown')
                    
                    return
                else:
                    bot.reply_to(message, LANG[lang]['wrong_file_type'])
                    return
            else:
                bot.reply_to(message, LANG[lang]['wrong_file_type'])
                return
        
        # Validate file type for selected service
        if not validate_file_for_service(message, service, lang):
            return
        
        if user_id in active_sessions:
            try:
                active_sessions[user_id]['timer'].cancel()
            except:
                pass
            active_sessions.pop(user_id, None)
            
        # Update user activity
        update_user_activity(user_id)
        
        # Show encrypting status
        lang = get_user_lang(message.from_user.language_code)
        status_msg = bot.send_message(message.chat.id, LANG[lang]['encrypting'])
        
        if message.content_type == 'photo':
            file_info = bot.get_file(message.photo[-1].file_id)
            original_name = f"photo_{message.photo[-1].file_id}.jpg"
        elif message.content_type == 'document':
            file_info = bot.get_file(message.document.file_id)
            original_name = message.document.file_name
        elif message.content_type == 'video':
            file_info = bot.get_file(message.video.file_id)
            original_name = f"video_{message.video.file_id}.mp4"
        elif message.content_type == 'audio':
            file_info = bot.get_file(message.audio.file_id)
            original_name = f"audio_{message.audio.file_id}.mp3"

        # Generate secure filename
        secure_filename = generate_secure_filename(original_name)
        
        # Download file
        downloaded_file = bot.download_file(file_info.file_path)
        file_path = f"files/{secure_filename}"
        
        # Get service for context
        service = user_services.get(user_id, 'general')
        
        # Get original file size for compression ratio calculation
        original_size = get_file_size_mb(downloaded_file)
        
        # Start encryption in a separate thread
        future = encryption_pool.submit(async_encrypt_file, downloaded_file, file_path)
        
        # Update encryption status with animated progress indicator
        animation_chars = ['⏳', '⌛', '⏳', '⌛']
        dots = 0
        char_idx = 0
        start_time = time.time()
        
        while not future.done():
            elapsed = time.time() - start_time
            dots = (dots % 3) + 1
            char_idx = (char_idx + 1) % len(animation_chars)
            
            try:
                # Show animated progress with elapsed time
                bot.edit_message_text(
                    f"{animation_chars[char_idx]} {LANG[lang]['encrypting']} {'.' * dots}\n({elapsed:.1f}s)",
                    message.chat.id, 
                    status_msg.message_id
                )
            except:
                pass
            
            # Shorter sleep time for smoother animation
            time.sleep(0.3)
        
        # Get encryption result
        try:
            encryption_time = future.result()
            # Show encryption completion message
            bot.edit_message_text(
                LANG[lang]['encryption_complete'].format(encryption_time),
                message.chat.id,
                status_msg.message_id
            )
        except Exception as e:
            logger.error(f"Encryption failed: {str(e)}")
            # If encryption in thread failed, try direct encryption as fallback
            encrypted_data, _ = encrypt_file(downloaded_file)
            with open(file_path, 'wb') as f:
                f.write(encrypted_data)

        # Store original name and secure path in database
        try:
            conn = sqlite3.connect('files.db', timeout=10.0)
            cursor = conn.execute('INSERT INTO files (user_id, file_id, file_name, file_path, created_at) VALUES (?, ?, ?, ?, ?)',
                        (message.from_user.id, file_info.file_id, original_name, file_path, datetime.now()))
            db_id = cursor.lastrowid
            conn.commit()
            conn.close()
        except Exception as db_error:
            logger.error(f"Database insert failed: {str(db_error)}")
            cleanup_failed_file(file_path)
            send_error_with_restart(message.chat.id, LANG[lang]['error_upload'], lang)
            return

        # Log secure handling
        logger.info(f"File securely stored: {original_name} → {secure_filename}")

        # Wait a moment to show the encryption success message
        time.sleep(1)
            
        file_type, ext = get_file_type(original_name)
        markup = types.InlineKeyboardMarkup()
        
        # Check if file format is supported
        if not is_supported_file(original_name):
            send_error_with_restart(message.chat.id, LANG[lang]['unsupported_format'], lang)
            cleanup_failed_file(file_path)
            return
        
        # Create contextual menu based on file type and service
        if file_type == 'image':
            options_text = LANG[lang]['image_options']
            if service == 'compress_image':
                markup.add(types.InlineKeyboardButton(LANG[lang]['compress_img'], callback_data=f"4_{db_id}"))
            else:
                if ext != 'jpg':
                    markup.add(types.InlineKeyboardButton(LANG[lang]['convert_jpg'], callback_data=f"1_{db_id}"))
                if ext != 'png':
                    markup.add(types.InlineKeyboardButton(LANG[lang]['convert_png'], callback_data=f"2_{db_id}"))
                if ext != 'webp':
                    markup.add(types.InlineKeyboardButton(LANG[lang]['convert_webp'], callback_data=f"3_{db_id}"))
                markup.add(types.InlineKeyboardButton(LANG[lang]['compress_img'], callback_data=f"4_{db_id}"))
        
        elif file_type == 'document':
            options_text = LANG[lang]['document_options']
            if service == 'pdf_merge':
                if HAS_PDF_MERGER:
                    markup.add(types.InlineKeyboardButton(LANG[lang]['combine_pdf'], callback_data=f"start_merge_{db_id}"))
            elif service == 'pdf_compress':
                markup.add(types.InlineKeyboardButton(LANG[lang]['compress_pdf'], callback_data=f"5_{db_id}"))
            elif service == 'pdf_convert':
                markup.add(types.InlineKeyboardButton('📄 Convert to Word', callback_data=f"8_{db_id}"))
            elif service == 'compress_zip':
                markup.add(types.InlineKeyboardButton(LANG[lang]['zip_file'], callback_data=f"7_{db_id}"))
            elif ext == 'pdf':
                markup.add(types.InlineKeyboardButton(LANG[lang]['compress_pdf'], callback_data=f"5_{db_id}"))
                markup.add(types.InlineKeyboardButton('📄 Convert to Word', callback_data=f"8_{db_id}"))
                if HAS_PDF_MERGER:
                    markup.add(types.InlineKeyboardButton(LANG[lang]['combine_pdf'], callback_data=f"start_merge_{db_id}"))
            else:
                markup.add(types.InlineKeyboardButton(LANG[lang]['convert_document'], callback_data=f"convert_pdf_{db_id}"))
                markup.add(types.InlineKeyboardButton(LANG[lang]['compress_document'], callback_data=f"compress_zip_{db_id}"))
        
        elif file_type == 'video':
            options_text = LANG[lang]['media_options']
            if ext != 'mp4':
                markup.add(types.InlineKeyboardButton(LANG[lang]['convert_to_mp4'], callback_data=f"video_mp4_{db_id}"))
            markup.add(types.InlineKeyboardButton(LANG[lang]['extract_mp3'], callback_data=f"6_{db_id}"))
            markup.add(types.InlineKeyboardButton(LANG[lang]['zip_file'], callback_data=f"7_{db_id}"))
        
        elif file_type == 'audio':
            options_text = LANG[lang]['media_options']
            if ext != 'mp3':
                markup.add(types.InlineKeyboardButton(LANG[lang]['convert_to_mp3'], callback_data=f"audio_mp3_{db_id}"))
            markup.add(types.InlineKeyboardButton(LANG[lang]['zip_file'], callback_data=f"7_{db_id}"))
        
        else:
            # For any other file type or compress_zip service
            options_text = 'What would you like to do with your file?'
            markup.add(types.InlineKeyboardButton(LANG[lang]['zip_file'], callback_data=f"7_{db_id}"))
        
        # Add back to menu button for all types
        markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
        markup.add(types.InlineKeyboardButton(LANG[lang]['cancel'], callback_data=f"cancel_{db_id}"))

        # Create contextual file message based on file type
        file_size_text = f" ({get_file_size_mb(downloaded_file):.1f} MB)" if get_file_size_mb(downloaded_file) >= 0.1 else ""
        file_message = f"✅ **{original_name}**{file_size_text}\n\n{options_text}\n\n{LANG[lang]['security_reminder']}"
        
        reply_msg = bot.reply_to(message, file_message, parse_mode='Markdown', reply_markup=markup)
        
        # Start the session timer
        start_session_timer(message.chat.id, file_path, db_id, lang)

    except Exception as e:
        logger.error(f"File handling error for user {message.from_user.id}: {str(e)}")
        lang = get_user_lang(message.from_user.language_code)
        
        # Clean up any partially created files
        try:
            if 'file_path' in locals():
                cleanup_failed_file(file_path)
        except:
            pass
            
        # Send user-friendly error message with restart button
        send_error_with_restart(message.chat.id, LANG[lang]['error_upload'], lang)

@bot.callback_query_handler(func=lambda call: True)
def callback_handler(call):
    """
    Handler utama untuk semua callback query (tombol inline keyboard).
    
    Parameter:
        call: Objek callback query dari Telegram
    
    Return:
        Tidak ada
    
    Catatan:
        - Menangani semua interaksi tombol dalam bot
        - Memproses berbagai aksi: konversi, kompresi, PDF merge, dll.
        - Melakukan validasi dan security check
        - Mengupdate aktivitas pengguna
        - Menangani file processing dengan status update
        - Membersihkan file temporary dan original setelah processing
        - Menampilkan hasil dengan opsi untuk processing lebih lanjut
        - Menangani error dengan cleanup yang aman
    """
    user_id = call.from_user.id
    username = call.from_user.username or 'Unknown'
    temp_files = []  # Track temporary files for cleanup
    
    try:
        logger.info(f"User {user_id} ({username}) clicked: {call.data}")
        lang = get_user_lang(call.from_user.language_code)
        
        # Update user activity when any button is clicked
        update_user_activity(user_id)
        
        # Handle start_bot callback for first welcome message
        if call.data == "start_bot" or call.data == "back_to_start":
            # Clear any service selection
            user_services.pop(user_id, None)
            
            # Clear any active PDF merge sessions
            if user_id in pdf_merge_sessions:
                clear_pdf_merge_session(user_id)
            
            # Clear any existing sessions to prevent duplication
            if user_id in active_sessions:
                try:
                    active_sessions[user_id]['timer'].cancel()
                    # Try to delete countdown message if exists
                    try:
                        bot.delete_message(call.message.chat.id, active_sessions[user_id]['countdown_msg_id'])
                    except:
                        pass
                except:
                    pass
                active_sessions.pop(user_id, None)
            
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton('📄 PDF Tools', callback_data="service_pdf"))
            markup.add(types.InlineKeyboardButton('📸 Image Tools', callback_data="service_image"))
            markup.add(types.InlineKeyboardButton('🎵 Media Tools', callback_data="service_media"))
            markup.add(types.InlineKeyboardButton('🗜️ Compression', callback_data="service_compress"))
            
            bot.edit_message_text(LANG[lang]['welcome'], call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
        
        # Handle service selection
        if call.data.startswith("service_"):
            service_type = call.data.split('_')[1]
            
            if service_type == 'pdf':
                markup = types.InlineKeyboardMarkup()
                if HAS_PDF_MERGER:
                    markup.add(types.InlineKeyboardButton('🔗 Merge Multiple PDFs', callback_data="pdf_merge"))
                markup.add(types.InlineKeyboardButton('🗜️ Compress PDF Size', callback_data="pdf_compress"))
                markup.add(types.InlineKeyboardButton('📄 PDF to Word Document', callback_data="pdf_convert"))
                markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
                
                bot.edit_message_text('📄 **PDF Tools**\n\nChoose what you want to do with your PDF:', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            
            elif service_type == 'image':
                user_services[user_id] = 'image'
                markup = types.InlineKeyboardMarkup()
                markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
                
                bot.edit_message_text('📸 **Image Tools**\n\nSend me any image file and I\'ll show you conversion and compression options.\n\n• Convert between JPG, PNG, WebP\n• Compress images to reduce size\n• Optimize image quality\n\n📱 **Ready for your image!**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            
            elif service_type == 'media':
                user_services[user_id] = 'media'
                markup = types.InlineKeyboardMarkup()
                markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
                
                bot.edit_message_text('🎵 **Media Tools**\n\nSend me any audio or video file for processing.\n\n• Convert video to MP4 format\n• Extract audio from videos\n• Convert audio to MP3\n• Compress media files\n\n🎬 **Ready for your media file!**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            
            elif service_type == 'compress':
                markup = types.InlineKeyboardMarkup()
                markup.add(types.InlineKeyboardButton('📄 Compress PDF Files', callback_data="pdf_compress"))
                markup.add(types.InlineKeyboardButton('📸 Compress Images', callback_data="compress_image"))
                markup.add(types.InlineKeyboardButton('📦 Create ZIP Archives', callback_data="compress_zip"))
                markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
                
                bot.edit_message_text('🗜️ **Compression Tools**\n\nReduce file sizes and create archives:\n\n• Compress PDF documents\n• Optimize image file sizes\n• Create ZIP archives for any files\n\n💾 **Choose compression type:**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            
            bot.answer_callback_query(call.id)
            return
        
        # Handle compression service selection
        if call.data == "compress_image":
            user_services[user_id] = 'compress_image'
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
            
            bot.edit_message_text('🗜️ **Image Compression**\n\nSend me any image file to reduce its size while maintaining quality.\n\n• Smart compression algorithms\n• Maintains visual quality\n• Reduces file size significantly\n\n📷 **Ready for your image!**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
        
        elif call.data == "compress_zip":
            user_services[user_id] = 'compress_zip'
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
            
            bot.edit_message_text('📦 **ZIP Archive Creation**\n\nSend me any file to create a compressed ZIP archive.\n\n• Universal file compression\n• Maximum compression level\n• Works with any file type\n• Easy to share and store\n\n📁 **Ready for your file!**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
        
        # Handle direct service selection (backward compatibility)
        if call.data == "pdf_merge":
            user_services[user_id] = 'pdf_merge'
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
            
            bot.edit_message_text('🔗 **PDF Merge Service**\n\nSend me 2 or more PDF files to combine them into one document.\n\n• Maintains original quality\n• Custom page ordering\n• Batch processing support\n• Secure file handling\n\n📄 **Ready for your PDFs!**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
        
        elif call.data == "pdf_split":
            user_services[user_id] = 'pdf_split'
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton(LANG[lang]['back_to_menu'], callback_data="back_to_start"))
            
            bot.edit_message_text("✂️ Upload a PDF file to split into separate pages", call.message.chat.id, call.message.message_id, reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
        
        elif call.data == "pdf_compress":
            user_services[user_id] = 'pdf_compress'
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
            
            bot.edit_message_text('🗜️ **PDF Compression**\n\nSend me a PDF file to reduce its size while maintaining readability.\n\n• Advanced compression algorithms\n• Maintains document quality\n• Reduces file size up to 70%\n• Perfect for sharing and storage\n\n📄 **Ready for your PDF!**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
        
        elif call.data == "pdf_convert":
            user_services[user_id] = 'pdf_convert'
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton('🔙 Back to Main Menu', callback_data="back_to_start"))
            
            bot.edit_message_text('📄 **PDF to Word Conversion**\n\nSend me a PDF file to convert it to an editable Word document.\n\n• Preserves formatting and layout\n• Maintains text and images\n• Creates editable DOCX files\n• Perfect for document editing\n\n📄 **Ready for your PDF!**', call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
        
        # Handle PDF merge callbacks
        if call.data.startswith("start_merge_"):
            if not HAS_PDF_MERGER:
                bot.answer_callback_query(call.id, "PDF merger not available")
                return
                
            db_id = call.data.split('_')[2]
            # Start PDF merge session
            create_pdf_merge_session(user_id, call.message.chat.id, lang)
            add_pdf_to_merge_session(user_id, db_id)
            user_services[user_id] = 'pdf_merge'
            
            bot.edit_message_text(
                "📄 **PDF Merge Started**\n\nSend me 2 or more PDF files to combine them.\n\n⏱️ **Collecting files...** (send multiple files now)\n\nI'll show you the order and let you rearrange before merging.",
                call.message.chat.id,
                call.message.message_id,
                parse_mode='Markdown'
            )
            
            bot.answer_callback_query(call.id, "PDF merge started - send more PDFs!")
            return
        
        if call.data == "show_reorder_options":
            # Show reordering interface
            session = get_pdf_merge_session(user_id)
            if session:
                pdf_list = generate_pdf_list_text(user_id, lang)
                markup = create_pdf_reorder_markup(user_id, lang)
                
                bot.edit_message_text(
                    f"✏️ **Reorder PDFs:**\n\n{pdf_list}\n\nUse ⬆️⬇️ to move files up/down, ❌ to remove files.",
                    call.message.chat.id,
                    call.message.message_id,
                    parse_mode='Markdown',
                    reply_markup=markup
                )
            
            bot.answer_callback_query(call.id)
            return
        
        if call.data == "back_to_confirmation":
            # Go back to order confirmation
            session = get_pdf_merge_session(user_id)
            if session:
                pdf_list = generate_pdf_list_text(user_id, lang)
                confirmation_text = f"📄 **Files ready ({len(session['pdfs'])} PDFs):**\n\n{pdf_list}\n\n✅ **Current merge order shown above.**\n\n❓ Do you want to change the order before merging?"
                
                markup = types.InlineKeyboardMarkup()
                markup.add(types.InlineKeyboardButton('🔗 Merge Now (Keep Order)', callback_data="execute_pdf_merge"))
                markup.add(types.InlineKeyboardButton('✏️ Change Order', callback_data="show_reorder_options"))
                markup.add(types.InlineKeyboardButton(LANG[lang]['cancel_merge'], callback_data="cancel_pdf_merge"))
                
                bot.edit_message_text(
                    confirmation_text,
                    call.message.chat.id,
                    call.message.message_id,
                    parse_mode='Markdown',
                    reply_markup=markup
                )
            
            bot.answer_callback_query(call.id)
            return
        
        if call.data == "execute_pdf_merge":
            # Execute PDF merge
            session = get_pdf_merge_session(user_id)
            if not session or len(session['pdfs']) < 2:
                bot.answer_callback_query(call.id, "Need at least 2 PDFs to merge")
                return
            
            status_msg = bot.send_message(call.message.chat.id, f"🔄 **Merging {len(session['pdfs'])} PDFs...**\n\nPlease wait while I combine your files.", parse_mode='Markdown')
            
            merged_data, error = merge_pdfs(user_id, lang)
            if merged_data:
                # Send merged PDF
                output = BytesIO(merged_data)
                bot.send_document(call.message.chat.id, output, visible_file_name="merged.pdf")
                
                # Calculate file size
                file_size_mb = get_file_size_mb(merged_data)
                bot.send_message(call.message.chat.id, f"✅ **PDFs merged successfully!**\n\n📄 Final file size: {file_size_mb:.1f} MB\n🔒 Original files deleted for security.", parse_mode='Markdown')
                
                # Clean up
                clear_pdf_merge_session(user_id)
                bot.delete_message(call.message.chat.id, status_msg.message_id)
                
                # Show completion options
                markup = types.InlineKeyboardMarkup()
                markup.add(
                    types.InlineKeyboardButton(LANG[lang]['yes_more'], callback_data="yes_more"),
                    types.InlineKeyboardButton(LANG[lang]['no_thanks'], callback_data="no_thanks")
                )
                bot.send_message(call.message.chat.id, LANG[lang]['help_more'], reply_markup=markup)
            else:
                bot.edit_message_text(f"❌ **Merge failed**\n\n{error}", call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                clear_pdf_merge_session(user_id)
            
            bot.answer_callback_query(call.id)
            return
        
        if call.data == "cancel_pdf_merge":
            session = get_pdf_merge_session(user_id)
            file_count = len(session['pdfs']) if session else 0
            
            clear_pdf_merge_session(user_id)
            bot.edit_message_text(
                f"❌ **PDF merge cancelled**\n\n{file_count} files deleted for security.",
                call.message.chat.id,
                call.message.message_id,
                parse_mode='Markdown'
            )
            bot.answer_callback_query(call.id, "Merge cancelled")
            return
        
        if call.data.startswith("move_pdf_up_"):
            index = int(call.data.split('_')[3])
            session = get_pdf_merge_session(user_id)
            if session and index > 0:
                # Swap with previous
                session['pdfs'][index], session['pdfs'][index-1] = session['pdfs'][index-1], session['pdfs'][index]
                
                # Update display
                pdf_list = generate_pdf_list_text(user_id, lang)
                markup = create_pdf_reorder_markup(user_id, lang)
                
                bot.edit_message_text(
                    f"✏️ **Reorder PDFs:**\n\n{pdf_list}\n\nUse ⬆️⬇️ to move files up/down, ❌ to remove files.",
                    call.message.chat.id,
                    call.message.message_id,
                    parse_mode='Markdown',
                    reply_markup=markup
                )
            
            bot.answer_callback_query(call.id, "Moved up ⬆️")
            return
        
        if call.data.startswith("move_pdf_down_"):
            index = int(call.data.split('_')[3])
            session = get_pdf_merge_session(user_id)
            if session and index < len(session['pdfs']) - 1:
                # Swap with next
                session['pdfs'][index], session['pdfs'][index+1] = session['pdfs'][index+1], session['pdfs'][index]
                
                # Update display
                pdf_list = generate_pdf_list_text(user_id, lang)
                markup = create_pdf_reorder_markup(user_id, lang)
                
                bot.edit_message_text(
                    f"✏️ **Reorder PDFs:**\n\n{pdf_list}\n\nUse ⬆️⬇️ to move files up/down, ❌ to remove files.",
                    call.message.chat.id,
                    call.message.message_id,
                    parse_mode='Markdown',
                    reply_markup=markup
                )
            
            bot.answer_callback_query(call.id, "Moved down ⬇️")
            return
        
        if call.data.startswith("remove_pdf_"):
            index = int(call.data.split('_')[2])
            session = get_pdf_merge_session(user_id)
            if session and 0 <= index < len(session['pdfs']):
                # Get filename for confirmation
                pdf_id = session['pdfs'][index]
                conn = sqlite3.connect('files.db')
                cursor = conn.execute('SELECT file_name FROM files WHERE id = ?', (pdf_id,))
                result = cursor.fetchone()
                filename = result[0] if result else "Unknown file"
                conn.close()
                
                # Remove PDF from session and clean up file
                session['pdfs'].pop(index)
                
                try:
                    conn = sqlite3.connect('files.db')
                    cursor = conn.execute('SELECT file_path FROM files WHERE id = ?', (pdf_id,))
                    result = cursor.fetchone()
                    if result:
                        cleanup_failed_file(result[0])
                        conn.execute('DELETE FROM files WHERE id = ?', (pdf_id,))
                        conn.commit()
                    conn.close()
                except Exception as e:
                    logger.error(f"Error removing PDF {pdf_id}: {str(e)}")
                
                # Update display or cancel if insufficient PDFs left
                if len(session['pdfs']) < 2:
                    clear_pdf_merge_session(user_id)
                    bot.edit_message_text(
                        f"❌ **Merge cancelled**\n\nNeed at least 2 PDFs to merge. Removed: {filename}",
                        call.message.chat.id,
                        call.message.message_id,
                        parse_mode='Markdown'
                    )
                else:
                    pdf_list = generate_pdf_list_text(user_id, lang)
                    markup = create_pdf_reorder_markup(user_id, lang)
                    
                    bot.edit_message_text(
                        f"✏️ **Reorder PDFs:**\n\n{pdf_list}\n\nUse ⬆️⬇️ to move files up/down, ❌ to remove files.\n\n✅ Removed: {filename}",
                        call.message.chat.id,
                        call.message.message_id,
                        parse_mode='Markdown',
                        reply_markup=markup
                    )
            
            bot.answer_callback_query(call.id, "File removed ❌")
            return
        
        if call.data == "noop":
            bot.answer_callback_query(call.id)
            return
        
        # Handle document conversion callbacks
        if call.data.startswith("convert_pdf_"):
            db_id = call.data.split('_')[2]
            # Redirect to Word to PDF conversion (action 9)
            call.data = f"9_{db_id}"
        
        elif call.data.startswith("compress_zip_"):
            db_id = call.data.split('_')[2]
            # Redirect to ZIP compression (action 7)
            call.data = f"7_{db_id}"
        
        # Handle video/audio conversion callbacks
        elif call.data.startswith("video_mp4_"):
            db_id = call.data.split('_')[2]
            # Redirect to video conversion (action 10)
            call.data = f"10_{db_id}"
        
        elif call.data.startswith("audio_mp3_"):
            db_id = call.data.split('_')[2]
            # Redirect to audio conversion (action 11)
            call.data = f"11_{db_id}"
        
        # Cancel session timer when user takes action
        if user_id in active_sessions:
            try:
                active_sessions[user_id]['timer'].cancel()
                # Try to delete the countdown message
                try:
                    bot.delete_message(call.message.chat.id, active_sessions[user_id]['countdown_msg_id'])
                except:
                    pass
            except:
                pass
            active_sessions.pop(user_id, None)
        
        if call.data == "start_over":
            bot.send_message(call.message.chat.id, LANG[lang]['ready_next'])
            bot.answer_callback_query(call.id)
            return
            
        if call.data == "yes_more":
            # Clear service selection and show comprehensive menu
            user_services.pop(user_id, None)
            
            # Clear any existing sessions to prevent duplication
            if user_id in active_sessions:
                try:
                    active_sessions[user_id]['timer'].cancel()
                    # Try to delete countdown message if exists
                    try:
                        bot.delete_message(call.message.chat.id, active_sessions[user_id]['countdown_msg_id'])
                    except:
                        pass
                except:
                    pass
                active_sessions.pop(user_id, None)
            
            # Clear any PDF merge sessions
            if user_id in pdf_merge_sessions:
                clear_pdf_merge_session(user_id)
            
            markup = types.InlineKeyboardMarkup()
            markup.add(types.InlineKeyboardButton('📄 PDF Tools', callback_data="service_pdf"))
            markup.add(types.InlineKeyboardButton('📸 Image Tools', callback_data="service_image"))
            markup.add(types.InlineKeyboardButton('🎵 Media Tools', callback_data="service_media"))
            markup.add(types.InlineKeyboardButton('🗜️ Compression', callback_data="service_compress"))
            
            bot.edit_message_text(LANG[lang]['welcome'], call.message.chat.id, call.message.message_id, parse_mode='Markdown', reply_markup=markup)
            bot.answer_callback_query(call.id)
            return
            
        if call.data == "no_thanks":
            bot.edit_message_text(LANG[lang]['goodbye'], 
                                call.message.chat.id, call.message.message_id)
            bot.answer_callback_query(call.id)
            return
            
        if '_' not in call.data or len(call.data.split('_')) < 2:
            bot.answer_callback_query(call.id, "❌ Invalid action!")
            return
        
        # Handle cancel action
        if call.data.startswith("cancel_"):
            db_id = call.data.split('_')[1]
            # Clean up file
            conn = sqlite3.connect('files.db')
            cursor = conn.execute('SELECT file_path FROM files WHERE id = ?', (db_id,))
            result = cursor.fetchone()
            if result:
                file_path = result[0]
                cleanup_failed_file(file_path)
                conn.execute('DELETE FROM files WHERE id = ?', (db_id,))
                conn.commit()
            conn.close()
            
            bot.edit_message_text("❌ Operation cancelled. File deleted for security.", 
                                call.message.chat.id, call.message.message_id)
            bot.answer_callback_query(call.id, "Cancelled")
            return
            
        try:
            action, db_id = call.data.split('_', 1)
        except ValueError:
            bot.answer_callback_query(call.id, "❌ Invalid action format!")
            return
        
        conn = sqlite3.connect('files.db')
        cursor = conn.execute('SELECT file_path, file_name FROM files WHERE id = ?', (db_id,))
        result = cursor.fetchone()
        conn.close()
        
        if not result:
            bot.answer_callback_query(call.id, "❌ File not found!")
            return
            
        file_path, original_name = result
        
        # Show processing status
        status_msg = bot.send_message(call.message.chat.id, LANG[lang]['compressing'])
        
        # Read and decrypt the file
        try:
            with open(file_path, 'rb') as f:
                encrypted_data = f.read()
            
            # Check if file is empty
            if not encrypted_data:
                raise ValueError("File is empty")
                
            file_data = decrypt_file(encrypted_data)
            
            # Get original file size for compression ratio calculation
            original_size = get_file_size_mb(file_data)
            
        except Exception as e:
            logger.error(f"Failed to read or decrypt file {file_path}: {str(e)}")
            
            # Clean up the corrupted file
            cleanup_failed_file(file_path)
            
            # Send user-friendly error with restart button
            send_error_with_restart(call.message.chat.id, LANG[lang]['error_processing'], lang)
            bot.answer_callback_query(call.id, "File processing error")
            return
        
        if action == "1":
            try:
                bot.edit_message_text('🔄 **Converting to JPG...**\n\nProcessing your image...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                img_io = BytesIO(file_data)
                with Image.open(img_io) as img:
                    # Handle different image modes properly
                    if img.mode in ('RGBA', 'LA', 'P'):
                        # Create white background for transparency
                        background = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                        img = background
                    elif img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    output = BytesIO()
                    img.save(output, format='JPEG', quality=95, optimize=True)
                    output.seek(0)
                    
                    converted_size = get_file_size_mb(output.getvalue())
                    bot.send_document(call.message.chat.id, output, visible_file_name="converted.jpg")
                    bot.send_message(call.message.chat.id, f'✅ **JPG conversion complete!**\n\n📄 File size: {converted_size:.1f} MB', parse_mode='Markdown')
            except Exception as e:
                logger.error(f"JPG conversion error: {str(e)}")
                bot.send_message(call.message.chat.id, f"❌ JPG conversion failed: {str(e)}")
        
        elif action == "2":
            try:
                bot.edit_message_text('🔄 **Converting to PNG...**\n\nProcessing your image...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                img_io = BytesIO(file_data)
                with Image.open(img_io) as img:
                    # Ensure proper PNG format
                    if img.mode not in ('RGBA', 'RGB', 'L', 'P'):
                        img = img.convert('RGBA')
                    
                    output = BytesIO()
                    img.save(output, format='PNG', optimize=True)
                    output.seek(0)
                    
                    converted_size = get_file_size_mb(output.getvalue())
                    bot.send_document(call.message.chat.id, output, visible_file_name="converted.png")
                    bot.send_message(call.message.chat.id, f'✅ **PNG conversion complete!**\n\n📄 File size: {converted_size:.1f} MB', parse_mode='Markdown')
            except Exception as e:
                logger.error(f"PNG conversion error: {str(e)}")
                bot.send_message(call.message.chat.id, f"❌ PNG conversion failed: {str(e)}")
        
        elif action == "3":
            try:
                bot.edit_message_text('🔄 **Converting to WebP...**\n\nProcessing your image...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                img_io = BytesIO(file_data)
                with Image.open(img_io) as img:
                    # WebP supports both RGB and RGBA
                    if img.mode not in ('RGB', 'RGBA'):
                        img = img.convert('RGB')
                    
                    output = BytesIO()
                    img.save(output, format='WEBP', quality=95, method=6)
                    output.seek(0)
                    
                    converted_size = get_file_size_mb(output.getvalue())
                    bot.send_document(call.message.chat.id, output, visible_file_name="converted.webp")
                    bot.send_message(call.message.chat.id, f'✅ **WebP conversion complete!**\n\n📄 File size: {converted_size:.1f} MB', parse_mode='Markdown')
            except Exception as e:
                logger.error(f"WebP conversion error: {str(e)}")
                bot.send_message(call.message.chat.id, f"❌ WebP conversion failed: {str(e)}")
        
        elif action == "4":
            try:
                bot.edit_message_text('🗜️ **Compressing image...**\n\nOptimizing file size...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                img_io = BytesIO(file_data)
                with Image.open(img_io) as img:
                    # Handle transparency properly
                    if img.mode in ('RGBA', 'LA', 'P'):
                        background = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                        img = background
                    elif img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    width, height = img.size
                    
                    # Smart compression based on image size
                    if width * height > 2000000:  # Large image (>2MP)
                        # Aggressive compression for large images
                        new_size = (int(width * 0.5), int(height * 0.5))
                        quality = 60
                    elif width * height > 500000:  # Medium image (>0.5MP)
                        # Moderate compression
                        new_size = (int(width * 0.7), int(height * 0.7))
                        quality = 70
                    else:
                        # Light compression for small images
                        new_size = (int(width * 0.8), int(height * 0.8))
                        quality = 80
                    
                    # Resize image
                    img_resized = img.resize(new_size, Image.Resampling.LANCZOS)
                    
                    output = BytesIO()
                    img_resized.save(output, format='JPEG', quality=quality, optimize=True)
                    output.seek(0)
                    
                    # Calculate compression ratio
                    compressed_size = get_file_size_mb(output.getvalue())
                    ratio = calculate_compression_ratio(original_size, compressed_size)
                    
                    # If still not compressed enough, try more aggressive settings
                    if ratio < MIN_COMPRESSION_TARGET and original_size > 1.0:  # Only for files > 1MB
                        output = BytesIO()
                        more_aggressive_size = (int(width * 0.4), int(height * 0.4))
                        img_more_compressed = img.resize(more_aggressive_size, Image.Resampling.LANCZOS)
                        img_more_compressed.save(output, format='JPEG', quality=50, optimize=True)
                        output.seek(0)
                        
                        # Recalculate compression ratio
                        new_compressed_size = get_file_size_mb(output.getvalue())
                        new_ratio = calculate_compression_ratio(original_size, new_compressed_size)
                        
                        # Use the better compression if it's significantly better
                        if new_ratio > ratio * 1.2:  # At least 20% better
                            compressed_size = new_compressed_size
                            ratio = new_ratio
                        else:
                            # Revert to previous compression
                            output = BytesIO()
                            img_resized.save(output, format='JPEG', quality=quality, optimize=True)
                            output.seek(0)
                    
                    # Send the compressed file
                    bot.send_document(call.message.chat.id, output, visible_file_name="compressed.jpg")
                    
                    # Show appropriate message based on compression ratio
                    if ratio < 0.1:  # Less than 10% compression
                        bot.send_message(call.message.chat.id, f'ℹ️ **Image already optimized**\n\nOriginal: {original_size:.1f} MB\nCompressed: {compressed_size:.1f} MB\n\nThis image is already well-optimized!', parse_mode='Markdown')
                    else:
                        # Show compression result with file sizes
                        savings = ((original_size - compressed_size) / original_size) * 100
                        bot.send_message(call.message.chat.id, f'✅ **Image compressed successfully!**\n\n📉 {original_size:.1f} MB → {compressed_size:.1f} MB\n💾 Space saved: {savings:.1f}%', parse_mode='Markdown')
                    
            except Exception as e:
                logger.error(f"Image compression error: {str(e)}")
                bot.send_message(call.message.chat.id, f"❌ Image compression failed: {str(e)}")
        
        elif action == "5":
            try:
                import fitz  # PyMuPDF
                pdf_io = BytesIO(file_data)
                doc = fitz.open(stream=pdf_io, filetype="pdf")
                output = BytesIO()
                
                new_doc = fitz.open()
                
                # More aggressive compression settings
                compression_matrix = fitz.Matrix(0.25, 0.25)  # 25% of original size
                jpg_quality = 10  # Lower quality for better compression
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    pix = page.get_pixmap(matrix=compression_matrix)
                    img_data = pix.tobytes("jpeg", jpg_quality=jpg_quality)
                    
                    new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                    new_page.insert_image(new_page.rect, stream=img_data)
                
                # Use maximum compression settings
                new_doc.save(output, garbage=4, deflate=True, clean=True, linear=True)
                new_doc.close()
                doc.close()
                output.seek(0)
                
                # Calculate compression ratio
                compressed_size = get_file_size_mb(output.getvalue())
                ratio = calculate_compression_ratio(original_size, compressed_size)
                
                # Check if we need better compression
                doc_pages = len(doc) if 'doc' in locals() else 1
                if ratio < MIN_COMPRESSION_TARGET and doc_pages > 1:
                    output = BytesIO()
                    doc = fitz.open(stream=pdf_io, filetype="pdf")
                    new_doc = fitz.open()
                    
                    # More balanced settings to maintain readability
                    compression_matrix = fitz.Matrix(0.3, 0.3)  # 30% of original size for better readability
                    jpg_quality = 20  # Better quality for readability
                    
                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        pix = page.get_pixmap(matrix=compression_matrix)
                        img_data = pix.tobytes("jpeg", jpg_quality=jpg_quality)
                        
                        new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
                        new_page.insert_image(new_page.rect, stream=img_data)
                    
                    new_doc.save(output, garbage=4, deflate=True, clean=True, linear=True)
                    new_doc.close()
                    doc.close()
                    output.seek(0)
                    
                    # Recalculate compression ratio
                    compressed_size = get_file_size_mb(output.getvalue())
                    ratio = calculate_compression_ratio(original_size, compressed_size)
                
                # Show appropriate message based on compression ratio
                if ratio < 0.1:  # Less than 10% compression
                    bot.send_message(call.message.chat.id, LANG[lang]['already_optimized'])
                else:
                    # Show compression result with file sizes
                    bot.send_message(call.message.chat.id, 
                                    LANG[lang]['compression_result'].format(original_size, compressed_size))
                
                # Send the compressed file
                bot.send_document(call.message.chat.id, output, visible_file_name="compressed.pdf")
                
                # Confirm file deletion for security
                bot.send_message(call.message.chat.id, LANG[lang]['files_deleted'])
            except Exception as fitz_error:
                logger.warning(f"PyMuPDF compression failed: {str(fitz_error)}")
                try:
                    temp_pdf = f"temp/temp_{db_id}_{int(time.time())}.pdf"
                    temp_files.append(temp_pdf)  # Track for cleanup
                    with open(temp_pdf, 'wb') as f:
                        f.write(file_data)
                    compressed_path = f"temp/compressed_{db_id}_{int(time.time())}.pdf"
                    temp_files.append(compressed_path)  # Track for cleanup
                    # Use more aggressive compression settings
                    subprocess.run(['gs', '-sDEVICE=pdfwrite', '-dPDFSETTINGS=/screen', 
                                  '-dDownsampleColorImages=true', '-dColorImageResolution=72',
                                  '-dCompatibilityLevel=1.4', '-dEmbedAllFonts=false',
                                  '-dSubsetFonts=true', '-dNOPAUSE', '-dQUIET', '-dBATCH', 
                                  f'-sOutputFile={compressed_path}', temp_pdf], 
                                 check=True, capture_output=True)
                    with open(compressed_path, 'rb') as f:
                        compressed_data = f.read()
                    
                    # Calculate compression ratio
                    compressed_size = get_file_size_mb(compressed_data)
                    ratio = calculate_compression_ratio(original_size, compressed_size)
                    
                    # If compression target not met but still preserving readability
                    if ratio < MIN_COMPRESSION_TARGET:
                        # Try with more balanced settings
                        more_compressed_path = f"temp/more_compressed_{db_id}_{int(time.time())}.pdf"
                        temp_files.append(more_compressed_path)
                        
                        # Use more balanced settings to maintain readability
                        subprocess.run(['gs', '-sDEVICE=pdfwrite', '-dPDFSETTINGS=/ebook', 
                                      '-dDownsampleColorImages=true', '-dColorImageResolution=150',
                                      '-dDownsampleGrayImages=true', '-dGrayImageResolution=150',
                                      '-dDownsampleMonoImages=true', '-dMonoImageResolution=150',
                                      '-dCompatibilityLevel=1.5', '-dEmbedAllFonts=true',
                                      '-dSubsetFonts=true', '-dNOPAUSE', '-dQUIET', '-dBATCH', 
                                      f'-sOutputFile={more_compressed_path}', temp_pdf], 
                                     check=True, capture_output=True)
                        
                        with open(more_compressed_path, 'rb') as f:
                            more_compressed_data = f.read()
                        
                        more_compressed_size = get_file_size_mb(more_compressed_data)
                        more_ratio = calculate_compression_ratio(original_size, more_compressed_size)
                        
                        # Use the better compression while maintaining readability
                        if more_ratio > ratio * 0.8:  # Accept if at least 80% as effective
                            compressed_data = more_compressed_data
                            compressed_size = more_compressed_size
                            ratio = more_ratio
                            os.remove(compressed_path)
                        else:
                            os.remove(more_compressed_path)
                    
                    # Show appropriate message based on compression ratio
                    if ratio < 0.1:  # Less than 10% compression
                        bot.send_message(call.message.chat.id, LANG[lang]['already_optimized'])
                    else:
                        # Show compression result with file sizes
                        bot.send_message(call.message.chat.id, 
                                        LANG[lang]['compression_result'].format(original_size, compressed_size))
                    
                    # Send as BytesIO to avoid file access issues
                    output = BytesIO(compressed_data)
                    output.seek(0)
                    bot.send_document(call.message.chat.id, output, visible_file_name="compressed.pdf")
                    
                    # Confirm file deletion for security
                    bot.send_message(call.message.chat.id, LANG[lang]['files_deleted'])
                    os.remove(compressed_path)
                    os.remove(temp_pdf)
                except Exception as e:
                    logger.error(f"PDF compression error: {str(e)}")
                    output = BytesIO(file_data)
                    output.seek(0)
                    bot.send_document(call.message.chat.id, output, visible_file_name="compressed.pdf")
        
        elif action == "6":
            try:
                bot.edit_message_text('🎵 **Extracting audio...**\n\nExtracting MP3 from video...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                output_path = f"temp/audio_{db_id}_{int(time.time())}.mp3"
                temp_files.append(output_path)  # Track for cleanup
                temp_video = f"temp/temp_video_{db_id}_{int(time.time())}"
                temp_files.append(temp_video)  # Track for cleanup
                
                with open(temp_video, 'wb') as f:
                    f.write(file_data)
                subprocess.run(['ffmpeg', '-i', temp_video, '-q:a', '0', '-map', 'a', output_path], 
                             check=True, capture_output=True)
                
                with open(output_path, 'rb') as f:
                    audio_data = f.read()
                    audio_size = get_file_size_mb(audio_data)
                    f.seek(0)
                    bot.send_audio(call.message.chat.id, f)
                
                bot.send_message(call.message.chat.id, f'✅ **Audio extracted successfully!**\n\n📄 File size: {audio_size:.1f} MB', parse_mode='Markdown')
                
                os.remove(output_path)
                os.remove(temp_video)
            except Exception as ffmpeg_error:
                logger.error(f"FFmpeg audio extraction failed: {str(ffmpeg_error)}")
                bot.answer_callback_query(call.id, "❌ Audio extraction failed!")
                bot.send_message(call.message.chat.id, '❌ **Audio extraction failed**\n\nFFmpeg may not be installed or the video has no audio track.', parse_mode='Markdown')
                return
        
        elif action == "7":
            try:
                bot.edit_message_text('📦 **Creating ZIP archive...**\n\nCompressing file...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                output = BytesIO()
                temp_file = f"temp/temp_zip_{db_id}_{int(time.time())}"
                temp_files.append(temp_file)  # Track for cleanup
                with open(temp_file, 'wb') as f:
                    f.write(file_data)
                    
                # Use maximum compression level (9) for better compression
                with zipfile.ZipFile(output, 'w', zipfile.ZIP_DEFLATED, compresslevel=9) as zipf:
                    zipf.write(temp_file, original_name)
                os.remove(temp_file)
                output.seek(0)
                
                # Calculate compression ratio
                compressed_size = get_file_size_mb(output.getvalue())
                ratio = calculate_compression_ratio(original_size, compressed_size)
                
                # Send the compressed file
                bot.send_document(call.message.chat.id, output, visible_file_name="compressed.zip")
                
                # Show appropriate message based on compression ratio
                if ratio < 0.1:  # Less than 10% compression
                    bot.send_message(call.message.chat.id, f'ℹ️ **File already compressed**\n\nOriginal: {original_size:.1f} MB\nZIP: {compressed_size:.1f} MB\n\nThis file type doesn\'t compress much further.', parse_mode='Markdown')
                else:
                    savings = ((original_size - compressed_size) / original_size) * 100
                    bot.send_message(call.message.chat.id, f'✅ **ZIP created successfully!**\n\n📉 {original_size:.1f} MB → {compressed_size:.1f} MB\n💾 Space saved: {savings:.1f}%', parse_mode='Markdown')
            except Exception as e:
                logger.error(f"ZIP compression error: {str(e)}")
                bot.send_message(call.message.chat.id, f"❌ ZIP compression failed: {str(e)}")
        
        # PDF to Word conversion
        elif action == "8":
            try:
                # Create a secure temporary directory that will be automatically cleaned up
                temp_dir = tempfile.mkdtemp(prefix="rupaganti_")
                try:
                    # Create temp PDF file with unique name - use normalized path
                    temp_pdf = os.path.normpath(os.path.join(temp_dir, f"temp_{db_id}.pdf"))
                    with open(temp_pdf, 'wb') as f:
                        f.write(file_data)
                    
                    # Output DOCX path - use normalized path
                    output_docx = os.path.normpath(os.path.join(temp_dir, f"converted_{db_id}.docx"))
                    
                    # Use pdf2docx for better conversion if available
                    pdf_conversion_success = False
                    if 'pdf2docx' in globals() or 'Converter' in globals():
                        try:
                            # Convert PDF to DOCX using pdf2docx
                            from pdf2docx import Converter
                            cv = Converter(temp_pdf)
                            cv.convert(output_docx, start=0, end=None)
                            cv.close()
                            pdf_conversion_success = True
                        except Exception as pdf_error:
                            logger.error(f"pdf2docx conversion failed: {str(pdf_error)}")
                            # Fall back to basic conversion
                            pdf_conversion_success = False
                    if not pdf_conversion_success:
                        # Fallback to basic conversion using python-docx and PyMuPDF
                        try:
                            from docx import Document
                            import fitz
                            
                            doc = Document()
                            with fitz.open(temp_pdf) as pdf_document:
                                for page_num in range(len(pdf_document)):
                                    page = pdf_document[page_num]
                                    # Get text with more formatting options
                                    text = page.get_text("text")
                                    if text.strip():  # Only add non-empty text
                                        doc.add_paragraph(text)
                                        
                                    # Try to extract images if text is limited
                                    if len(text.strip()) < 100:  # Likely image-heavy page
                                        try:
                                            # Add a note about possible image content
                                            doc.add_paragraph("[This page may contain images that couldn't be converted to text]")
                                        except:
                                            pass
                            
                            doc.save(output_docx)
                        except Exception as basic_error:
                            logger.error(f"Basic PDF conversion failed: {str(basic_error)}")
                            # Create an empty document with error message
                            try:
                                from docx import Document
                                doc = Document()
                                doc.add_paragraph("Error converting PDF. The file may be encrypted or contain only images.")
                                doc.save(output_docx)
                            except Exception:
                                # If even this fails, create a simple text file
                                with open(output_docx, 'w') as f:
                                    f.write("Error converting PDF. The file may be encrypted or contain only images.")
                    
                    # Make sure to close all file handles before sending
                    with open(output_docx, 'rb') as f:
                        file_content = f.read()
                    
                    # Send as BytesIO to avoid file access issues
                    output = BytesIO(file_content)
                    output.name = "converted.docx"
                    bot.send_document(call.message.chat.id, output, visible_file_name="converted.docx")
                    
                    # Clean up temporary directory manually
                finally:
                    try:
                        shutil.rmtree(temp_dir)
                    except Exception as cleanup_error:
                        logger.error(f"Failed to clean up temp directory: {str(cleanup_error)}")
                
                # Send a message about the conversion quality
                if pdf_conversion_success:
                    bot.send_message(call.message.chat.id, "✅ PDF converted to Word using enhanced conversion engine")
                else:
                    bot.send_message(call.message.chat.id, "✅ PDF converted to Word (basic conversion)")
                    
                # Confirm file deletion for security
                bot.send_message(call.message.chat.id, LANG[lang]['files_deleted'])
                    
                # Securely delete the original file immediately
                if os.path.exists(file_path):
                    try:
                        os.remove(file_path)
                        logger.info(f"Original file deleted after processing: {file_path}")
                    except Exception as e:
                        logger.error(f"Failed to delete original file {file_path}: {str(e)}")
                
            except Exception as e:
                logger.error(f"PDF to Word conversion error: {str(e)}")
                send_error_with_restart(call.message.chat.id, f"❌ PDF to Word conversion failed. {LANG[lang]['try_again']}", lang)
        
        # Video to MP4 conversion
        elif action == "10":
            try:
                bot.edit_message_text('🎬 **Converting to MP4...**\n\nThis may take a moment for large videos...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                temp_input = f"temp/input_{db_id}_{int(time.time())}"
                temp_output = f"temp/output_{db_id}_{int(time.time())}.mp4"
                temp_files.extend([temp_input, temp_output])
                
                with open(temp_input, 'wb') as f:
                    f.write(file_data)
                
                # Convert to MP4 using ffmpeg
                subprocess.run(['ffmpeg', '-i', temp_input, '-c:v', 'libx264', '-c:a', 'aac', 
                               '-preset', 'fast', '-crf', '23', temp_output], 
                              check=True, capture_output=True)
                
                with open(temp_output, 'rb') as f:
                    converted_data = f.read()
                
                converted_size = get_file_size_mb(converted_data)
                output = BytesIO(converted_data)
                bot.send_document(call.message.chat.id, output, visible_file_name="converted.mp4")
                bot.send_message(call.message.chat.id, f'✅ **MP4 conversion complete!**\n\n📄 File size: {converted_size:.1f} MB', parse_mode='Markdown')
                
                # Clean up temp files
                for temp_file in [temp_input, temp_output]:
                    try:
                        if os.path.exists(temp_file):
                            os.remove(temp_file)
                    except Exception as cleanup_error:
                        logger.error(f"Failed to cleanup temp file {temp_file}: {str(cleanup_error)}")
                        
            except Exception as e:
                logger.error(f"Video conversion error: {str(e)}")
                bot.send_message(call.message.chat.id, f"❌ Video conversion failed. FFmpeg may not be installed.")
        
        # Audio to MP3 conversion
        elif action == "11":
            try:
                bot.edit_message_text('🎵 **Converting to MP3...**\n\nProcessing audio...', call.message.chat.id, status_msg.message_id, parse_mode='Markdown')
                
                temp_input = f"temp/input_{db_id}_{int(time.time())}"
                temp_output = f"temp/output_{db_id}_{int(time.time())}.mp3"
                temp_files.extend([temp_input, temp_output])
                
                with open(temp_input, 'wb') as f:
                    f.write(file_data)
                
                # Convert to MP3 using ffmpeg
                subprocess.run(['ffmpeg', '-i', temp_input, '-c:a', 'libmp3lame', 
                               '-b:a', '192k', temp_output], 
                              check=True, capture_output=True)
                
                with open(temp_output, 'rb') as f:
                    converted_data = f.read()
                
                converted_size = get_file_size_mb(converted_data)
                output = BytesIO(converted_data)
                bot.send_audio(call.message.chat.id, output, title="Converted Audio")
                bot.send_message(call.message.chat.id, f'✅ **MP3 conversion complete!**\n\n📄 File size: {converted_size:.1f} MB', parse_mode='Markdown')
                
                # Clean up temp files
                for temp_file in [temp_input, temp_output]:
                    try:
                        if os.path.exists(temp_file):
                            os.remove(temp_file)
                    except Exception as cleanup_error:
                        logger.error(f"Failed to cleanup temp file {temp_file}: {str(cleanup_error)}")
                        
            except Exception as e:
                logger.error(f"Audio conversion error: {str(e)}")
                bot.send_message(call.message.chat.id, f"❌ Audio conversion failed. FFmpeg may not be installed.")
        
        # Word to PDF conversion
        elif action == "9":
            try:
                bot.edit_message_text(LANG[lang]['converting_to_pdf'], call.message.chat.id, status_msg.message_id)
            except:
                pass
                
            try:
                temp_dir = tempfile.mkdtemp(prefix="rupaganti_")
                temp_docx = os.path.join(temp_dir, f"input_{db_id}.docx")
                output_pdf = os.path.join(temp_dir, f"output_{db_id}.pdf")
                
                try:
                    # Save DOCX file
                    with open(temp_docx, 'wb') as f:
                        f.write(file_data)
                    
                    conversion_success = False
                    
                    # Try docx2pdf first
                    try:
                        import importlib.util
                        if importlib.util.find_spec("docx2pdf"):
                            from docx2pdf import convert
                            convert(temp_docx, output_pdf)
                            if os.path.exists(output_pdf) and os.path.getsize(output_pdf) > 0:
                                conversion_success = True
                    except Exception as e:
                        logger.warning(f"docx2pdf failed: {str(e)}")
                    
                    # Fallback to basic conversion for all document types
                    if not conversion_success:
                        try:
                            file_type, ext = get_file_type(original_name)
                            if ext in ['docx', 'doc']:
                                from docx import Document
                                doc = Document(temp_docx)
                                text_content = [para.text for para in doc.paragraphs if para.text.strip()]
                            elif ext in ['txt', 'rtf']:
                                with open(temp_docx, 'r', encoding='utf-8', errors='ignore') as f:
                                    text_content = f.readlines()
                            elif ext in ['xlsx', 'xls']:
                                try:
                                    import pandas as pd
                                    df = pd.read_excel(temp_docx)
                                    text_content = [df.to_string()]
                                except:
                                    text_content = ["Excel file content (conversion limited)"]
                            elif ext in ['pptx', 'ppt']:
                                try:
                                    from pptx import Presentation
                                    prs = Presentation(temp_docx)
                                    text_content = []
                                    for slide in prs.slides:
                                        for shape in slide.shapes:
                                            if hasattr(shape, "text"):
                                                text_content.append(shape.text)
                                except:
                                    text_content = ["PowerPoint file content (conversion limited)"]
                            else:
                                text_content = ["Document content"]
                            
                            # Create PDF with extracted content
                            from reportlab.pdfgen import canvas
                            from reportlab.lib.pagesizes import letter
                            
                            c = canvas.Canvas(output_pdf, pagesize=letter)
                            width, height = letter
                            y = height - 50
                            
                            for line in text_content:
                                if line and line.strip():
                                    text = line.strip()[:80]  # Limit line length
                                    c.drawString(50, y, text)
                                    y -= 20
                                    if y < 50:
                                        c.showPage()
                                        y = height - 50
                            
                            c.save()
                            conversion_success = True
                        except Exception as e:
                            logger.error(f"Fallback conversion failed: {str(e)}")
                    
                    if not conversion_success:
                        raise Exception("Conversion failed")
                    
                    # Send PDF
                    with open(output_pdf, 'rb') as f:
                        file_content = f.read()
                    
                    output = BytesIO(file_content)
                    filename = original_name.rsplit('.', 1)[0] + '.pdf'
                    bot.send_document(call.message.chat.id, output, visible_file_name=filename)
                    
                    bot.send_message(call.message.chat.id, LANG[lang]['pdf_conversion_success'])
                    bot.send_message(call.message.chat.id, LANG[lang]['file_ready'])
                    
                    file_size_mb = get_file_size_mb(file_content)
                    bot.send_message(call.message.chat.id, f"📄 PDF created ({file_size_mb:.1f} MB)")
                    
                finally:
                    try:
                        shutil.rmtree(temp_dir)
                    except Exception as cleanup_error:
                        logger.error(f"Failed to cleanup temp directory: {str(cleanup_error)}")
                
                cleanup_failed_file(file_path)
                bot.send_message(call.message.chat.id, LANG[lang]['files_deleted'])
                
            except Exception as e:
                logger.error(f"Word to PDF error: {str(e)}")
                cleanup_failed_file(file_path)
                send_error_with_restart(call.message.chat.id, LANG[lang]['pdf_conversion_failed'], lang)
        
        # Clean up original file after processing
        if action not in ["8", "9"]:
            try:
                conn = sqlite3.connect('files.db')
                conn.execute('DELETE FROM files WHERE id = ?', (db_id,))
                conn.commit()
                conn.close()
                
                if os.path.exists(file_path):
                    os.remove(file_path)
                    logger.info(f"Original file deleted after processing: {file_path}")
            except Exception as e:
                logger.error(f"Failed to delete original file {file_path}: {str(e)}")
        
        # Update status message to cleaning
        try:
            bot.edit_message_text(LANG[lang]['cleaning'], call.message.chat.id, status_msg.message_id)
        except:
            pass
            
        # Send completion message with Yes/No buttons
        markup = types.InlineKeyboardMarkup()
        markup.add(
            types.InlineKeyboardButton(LANG[lang]['yes_more'], callback_data="yes_more"),
            types.InlineKeyboardButton(LANG[lang]['no_thanks'], callback_data="no_thanks")
        )
        
        bot.send_message(call.message.chat.id, LANG[lang]['complete'])
        bot.send_message(call.message.chat.id, LANG[lang]['help_more'], reply_markup=markup)
        bot.answer_callback_query(call.id, LANG[lang]['done'])
        
        # Delete status message
        try:
            if 'status_msg' in locals():
                bot.delete_message(call.message.chat.id, status_msg.message_id)
        except Exception as delete_error:
            logger.debug(f"Could not delete status message: {str(delete_error)}")
        
    except Exception as e:
        logger.error(f"Callback handler error for user {user_id}: {str(e)}", exc_info=True)
        logger.error(f"Error details: {type(e).__name__}: {str(e)}")
        
        # Clean up any temporary files that might have been created
        for temp_file in temp_files:
            try:
                if temp_file and os.path.exists(temp_file):
                    os.remove(temp_file)
                    logger.info(f"Cleaned up temp file after error: {temp_file}")
            except Exception as cleanup_error:
                logger.error(f"Failed to clean up temp file {temp_file}: {str(cleanup_error)}")
        
        # Clean up original file if it exists
        try:
            if 'file_path' in locals() and file_path:
                cleanup_failed_file(file_path)
                if 'db_id' in locals():
                    conn = sqlite3.connect('files.db')
                    conn.execute('DELETE FROM files WHERE id = ?', (db_id,))
                    conn.commit()
                    conn.close()
        except Exception as cleanup_error:
            logger.error(f"Failed to cleanup after error: {str(cleanup_error)}")
                
        # Send user-friendly error message with restart button
        try:
            bot.answer_callback_query(call.id, "Processing error")
            send_error_with_restart(call.message.chat.id, LANG[lang]['oops_error'], lang)
        except Exception as error_send_error:
            logger.error(f"Failed to send error message: {str(error_send_error)}")


if __name__ == "__main__":
    """
    Entry point utama aplikasi bot.
    
    Fungsi:
        - Menampilkan pesan startup dengan jenis enkripsi yang digunakan
        - Membersihkan file dan database dari run sebelumnya
        - Memvalidasi format bot token
        - Memulai bot polling dengan pengaturan yang robust
        - Menangani berbagai jenis error dengan pesan yang informatif
    
    Catatan:
        - Menggunakan AES-256 jika tersedia, fallback ke Fernet
        - Membersihkan direktori 'files' dan 'temp' saat startup
        - Polling dengan interval 1 detik dan timeout 20 detik
        - Memberikan pesan error yang spesifik untuk troubleshooting
        - Exit dengan kode 1 jika terjadi error kritis
    """
    # Security startup message
    encryption_type = "AES-256 with hardware acceleration" if HAS_AES else "Fernet"
    print(f"🚀 Bot started securely with {encryption_type}... waiting for file uploads 🛡️")
    logger.info(f"Secure RupaGanti Bot starting with enhanced {encryption_type} encryption...")
    
    # Clean any leftover files from previous runs
    try:
        # First clean the database to remove references to files that might not exist
        try:
            conn = sqlite3.connect('files.db')
            conn.execute('DELETE FROM files')
            conn.commit()
            conn.close()
            logger.info("Database cleaned on startup")
        except Exception as db_error:
            logger.error(f"Database cleanup error: {str(db_error)}")
        
        # Then clean all files in storage directories
        for folder in ["files", "temp"]:
            if os.path.exists(folder):
                for filename in os.listdir(folder):
                    file_path = os.path.join(folder, filename)
                    if os.path.isfile(file_path):
                        try:
                            os.remove(file_path)
                            logger.info(f"Startup cleanup: Deleted {file_path}")
                        except Exception as file_error:
                            logger.error(f"Could not delete file {file_path}: {str(file_error)}")
        logger.info("Initial cleanup completed - all previous files deleted")
    except Exception as e:
        logger.error(f"Initial cleanup error: {str(e)}")
    
    # Verify bot token format
    if not BOT_TOKEN or len(BOT_TOKEN.split(':')) != 2:
        logger.critical("Invalid bot token format. Please check your token.")
        print("ERROR: Invalid bot token format. Please check your token.")
        exit(1)
        
    try:
        # Use more robust polling settings
        bot.polling(none_stop=True, interval=1, timeout=20)
    except Exception as e:
        logger.critical(f"Bot crashed: {str(e)}", exc_info=True)
        print(f"ERROR: {str(e)}")
        
        # Provide helpful error message for common issues
        if "401" in str(e):
            print("\nAuthentication failed: Your bot token is invalid or revoked.")
            print("Please get a new token from @BotFather on Telegram.")
        elif "404" in str(e):
            print("\nAPI endpoint not found: Check your internet connection.")
            print("If you're using a proxy, make sure it's configured correctly.")
        elif "429" in str(e):
            print("\nToo many requests: You're being rate limited by Telegram.")
            print("Wait a while before trying again.")
        
        print("\nFor help, visit: https://core.telegram.org/bots/api")
        exit(1)